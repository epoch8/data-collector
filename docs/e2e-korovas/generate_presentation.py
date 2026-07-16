"""
Generate E2E Korovas pipeline presentation.

Run: python docs/e2e-korovas/generate_presentation.py
"""

from __future__ import annotations

import importlib.util
import sys
from pathlib import Path

from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent
REPO = ROOT.parent.parent
BUSINESS = ROOT.parent / "business"
E2E_IMG = ROOT / "img"
DATAPIPE_IMG = E2E_IMG / "datapipe"
VIDEO = ROOT / "video"
POSTERS = VIDEO / "posters"
LOGO = REPO / "e8-team-logo-1024.png"
OUT = ROOT / "Korovas-E2E.pptx"

spec = importlib.util.spec_from_file_location("gen_training", BUSINESS / "generate_training_presentation.py")
gen = importlib.util.module_from_spec(spec)
sys.modules["gen_training"] = gen
spec.loader.exec_module(gen)

C = gen.C
W = gen.W
TOTAL = 18

bg = gen.bg
title = gen.title
text = gen.text
rect = gen.rect
oval = gen.oval
pill = gen.pill
card = gen.card
arrow = gen.arrow
add_logo = gen.add_logo
new_deck = gen.new_deck


def footer_e2e(slide, n: int, dark: bool = False) -> None:
    from pptx.dml.color import RGBColor
    color = RGBColor(0x7C, 0x83, 0x90) if dark else C.faint
    text(slide, Inches(0.65), Inches(7.1), Inches(5.0), Inches(0.22), "Epoch8 · Korovas E2E", size=8, color=color)
    text(slide, Inches(12.15), Inches(7.1), Inches(0.7), Inches(0.22), f"{n:02d}/{TOTAL:02d}", size=8, color=color, align=PP_ALIGN.RIGHT)


def _fit_size(iw: int, ih: int, max_w, max_h) -> tuple[int, int]:
    if iw <= 0 or ih <= 0:
        return int(max_w), int(max_h)
    scale = min(float(max_w) / iw, float(max_h) / ih)
    return int(iw * scale), int(ih * scale)


def _resolve_media(folder: Path) -> Path | None:
    if not folder.is_dir():
        return None
    for name in ("screenshot.png", "screenshot.gif", "screenshot.jpg", "screenshot.webp"):
        p = folder / name
        if p.exists():
            return p
    for ext in (".png", ".gif", ".jpg", ".jpeg", ".webp"):
        for p in sorted(folder.glob(f"*{ext}")):
            if p.name != ".gitkeep":
                return p
    return None


def add_media(slide, path: Path | None, l, t, max_w, max_h, *, caption: str | None = None, accent=C.cobalt, frame: bool = True, theme: str = "light", placeholder: str | None = None):
    from pptx.dml.color import RGBColor
    fill = C.night_2 if theme == "dark" else C.white
    border = C.line if theme == "light" else RGBColor(0x35, 0x3D, 0x4C)

    if path is None or not path.exists():
        rect(slide, l, t, max_w, max_h, C.paper_2, border, radius=True)
        label = placeholder or "Скриншот — добавить позже"
        text(slide, l + Inches(0.2), t + max_h / 2 - Inches(0.35), max_w - Inches(0.4), Inches(0.7), label, size=10, color=C.muted, align=PP_ALIGN.CENTER)
        if caption:
            gen.add_caption(slide, l, t + max_h - Inches(0.18), max_w, caption, accent, theme=theme)
        return None

    iw, ih = 1, 1
    if path.suffix.lower() != ".gif":
        from PIL import Image
        with Image.open(path) as img:
            iw, ih = img.size

    final_w, final_h = _fit_size(iw, ih, max_w, max_h)
    final_l = l + int((max_w - final_w) / 2)
    final_t = t + int((max_h - final_h) / 2)

    if frame:
        rect(slide, l, t, max_w, max_h, fill, border, radius=True)

    pic = slide.shapes.add_picture(str(path), final_l, final_t, width=final_w, height=final_h)
    if caption:
        gen.add_caption(slide, l, t + max_h - Inches(0.18), max_w, caption, accent, theme=theme)
    return pic


def add_video(slide, movie: Path, poster: Path | None, l, t, w, h):
    """Embed mp4; falls back to poster/placeholder if file missing."""
    from pptx.dml.color import RGBColor

    border = RGBColor(0x35, 0x3D, 0x4C)
    if not movie.exists():
        rect(slide, l, t, w, h, C.night_2, border, radius=True)
        text(slide, l, t + h / 2 - Inches(0.2), w, Inches(0.4), "▶  Видео — добавить файл", size=14, bold=True, color=C.white, align=PP_ALIGN.CENTER)
        return None
    poster_path = str(poster) if poster and poster.exists() else None
    return slide.shapes.add_movie(str(movie), l, t, w, h, poster_frame_image=poster_path, mime_type="video/mp4")


def _fit_box(iw: int, ih: int, max_w, max_h) -> tuple:
    """Return (w, h) fitting inside max box, preserving aspect ratio."""
    scale = min(float(max_w) / iw, float(max_h) / ih)
    return int(iw * scale), int(ih * scale)


def slide_video_placeholder(prs, n: int, headline: str, subtitle: str = "Видео-каст — вставить запись"):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, C.night)
    title(s, headline, subtitle, dark=True)
    from pptx.dml.color import RGBColor
    rect(s, Inches(1.5), Inches(2.2), Inches(10.3), Inches(4.2), C.night_2, RGBColor(0x35, 0x3D, 0x4C), radius=True)
    text(s, Inches(1.5), Inches(3.85), Inches(10.3), Inches(0.5), "▶  Место для видео", size=22, bold=True, color=C.white, align=PP_ALIGN.CENTER)
    footer_e2e(s, n, dark=True)


def title_e2e(slide, title_text: str, subtitle: str | None = None, *, dark: bool = False, subtitle_top=None):
    """Same as shared title, but allows a tighter subtitle offset (manual pptx tweaks)."""
    from pptx.dml.color import RGBColor
    color = C.white if dark else C.ink
    muted = RGBColor(0xB0, 0xB7, 0xC5) if dark else C.muted
    text(slide, Inches(0.65), Inches(0.55), Inches(10.2), Inches(0.9), title_text, size=30, bold=True, color=color)
    rect(slide, Inches(0.65), Inches(1.48), Inches(1.8), Pt(5), C.lime if dark else C.cobalt)
    if subtitle:
        st = subtitle_top if subtitle_top is not None else Inches(1.65)
        text(slide, Inches(0.65), st, Inches(10.2), Inches(0.45), subtitle, size=12, color=muted)


def slide_video_desktop(
    prs,
    n: int,
    headline: str,
    subtitle: str,
    movie: Path,
    poster: Path | None,
    *,
    vw: int = 1920,
    vh: int = 1080,
    layout: dict | None = None,
):
    """Landscape video-cast with optional locked layout from manual pptx edits."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, C.night)

    from pptx.dml.color import RGBColor
    muted = RGBColor(0xAE, 0xB7, 0xC7)
    border = RGBColor(0x2D, 0x36, 0x48)

    # Locked layouts captured from manual PowerPoint tweaks (do not auto-refit).
    if layout:
        title_e2e(s, headline, subtitle, dark=True, subtitle_top=Inches(layout.get("subtitle_top", 1.59)))
        frame = layout["frame"]  # left, top, width, height
        video = layout["video"]
        rect(s, Inches(frame[0]), Inches(frame[1]), Inches(frame[2]), Inches(frame[3]), C.night_2, border, radius=True)
        add_video(s, movie, poster, Inches(video[0]), Inches(video[1]), Inches(video[2]), Inches(video[3]))
    else:
        title_e2e(s, headline, subtitle, dark=True)
        area_l, area_t = Inches(0.75), Inches(1.95)
        area_w, area_h = Inches(11.85), Inches(4.85)
        vid_w, vid_h = _fit_box(vw, vh, area_w, area_h)
        vid_l = area_l + int((area_w - vid_w) / 2)
        vid_t = area_t + int((area_h - vid_h) / 2)
        pad = Inches(0.08)
        rect(s, vid_l - pad, vid_t - pad, vid_w + pad * 2, vid_h + pad * 2, C.night_2, border, radius=True)
        add_video(s, movie, poster, vid_l, vid_t, vid_w, vid_h)

    text(s, Inches(0.85), Inches(6.86), Inches(11.65), Inches(0.2), "Кликните по видео для воспроизведения", size=8, color=muted, align=PP_ALIGN.CENTER)
    footer_e2e(s, n, dark=True)


# Manual layouts from user-edited Korovas-E2E.pptx (captured 2026-07-16).
LAYOUT_VIDEO_FORM = {
    "subtitle_top": 1.591,
    "frame": (2.284, 1.870, 8.782, 5.010),
    "video": (2.669, 2.109, 8.069, 4.539),
}
LAYOUT_VIDEO_VIZ = {
    "subtitle_top": 1.589,
    "frame": (1.678, 1.870, 9.993, 5.010),
    "video": (2.218, 2.161, 8.897, 4.388),
}


def slide_video_phone(prs, n: int, headline: str, subtitle: str, movie: Path, poster: Path | None):
    """Portrait phone video with steps on the left — fills landscape slide without looking empty."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, C.night)
    title_e2e(s, headline, subtitle, dark=True)

    from pptx.dml.color import RGBColor
    muted = RGBColor(0xAE, 0xB7, 0xC7)
    border = RGBColor(0x2D, 0x36, 0x48)

    steps = [
        ("1", "Открыть форму", "Выбрать сценарий сбора в приложении", C.cobalt),
        ("2", "Пройти шаги", "Поля, ракурсы, проверка качества кадра", C.coral),
        ("3", "Сохранить пакет", "Пакет готов к загрузке на сервер", C.mint),
    ]
    y = Inches(2.0)
    for num, head, body, color in steps:
        rect(s, Inches(0.75), y, Inches(5.35), Inches(1.05), C.night_2, border, radius=True)
        rect(s, Inches(0.75), y, Inches(0.08), Inches(1.05), color)
        pill(s, Inches(0.95), y + Inches(0.18), Inches(0.5), num, color, C.white)
        text(s, Inches(1.6), y + Inches(0.16), Inches(4.2), Inches(0.3), head, size=14, bold=True, color=C.white)
        text(s, Inches(1.6), y + Inches(0.5), Inches(4.2), Inches(0.4), body, size=10, color=muted)
        y += Inches(1.2)

    text(s, Inches(0.75), Inches(5.75), Inches(5.35), Inches(0.5), "Видео снято с экрана телефона — сценарий целиком", size=10, color=muted)

    # phone bezel sized to 574×1280 (~9:20)
    phone_w, phone_h = Inches(2.55), Inches(5.25)
    phone_l = Inches(8.85)
    phone_t = Inches(1.7)
    rect(s, phone_l, phone_t, phone_w, phone_h, RGBColor(0x0A, 0x0E, 0x16), border, radius=True)
    oval(s, phone_l + Inches(0.75), phone_t + Inches(0.12), Inches(1.05), Inches(0.16), RGBColor(0x1A, 0x22, 0x32))
    pad = Inches(0.14)
    inner_w = phone_w - pad * 2
    inner_h = phone_h - Inches(0.52)
    vid_w, vid_h = _fit_box(574, 1280, inner_w, inner_h)
    vid_l = phone_l + pad + int((inner_w - vid_w) / 2)
    vid_t = phone_t + Inches(0.36) + int((inner_h - vid_h) / 2)
    add_video(s, movie, poster, vid_l, vid_t, vid_w, vid_h)
    footer_e2e(s, n, dark=True)


def slide_users(prs, n: int):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, C.night)
    # subtitle_top locked from manual pptx edit
    title_e2e(s, "Пользователи и права доступа", "Создание пользователя и выдача доступа к проектам", dark=True, subtitle_top=Inches(1.538))

    from pptx.dml.color import RGBColor
    from PIL import Image

    muted = RGBColor(0xAE, 0xB7, 0xC7)
    border = RGBColor(0x2D, 0x36, 0x48)

    user_list = E2E_IMG / "user" / "user-list-crop.png"
    if not user_list.exists():
        user_list = E2E_IMG / "user" / "user list.png"
    access = E2E_IMG / "user" / "access-list-crop.png"
    if not access.exists():
        access = E2E_IMG / "user" / "accses_list.png"

    # Layout locked from manual pptx edit (2026-07-16):
    # two equal panels, screenshots width-filled, callouts pinned at bottom.
    panels = [
        (
            Inches(0.75), "1", "Выбрать пользователя",
            "Синхронизировать список с Firebase и открыть «Настроить»",
            user_list, C.cobalt, "Список пользователей", "Настроить →",
        ),
        (
            Inches(6.75), "2", "Выдать права",
            "Отметить проекты для мобилки и client-admin, затем сохранить",
            access, C.mint, "Права доступа", "Мобилка · Client-admin → Сохранить",
        ),
    ]
    panel_t, panel_w, panel_h = Inches(1.85), Inches(5.8), Inches(4.9)
    for left, num, head, body, path, accent, caption, callout in panels:
        rect(s, left, panel_t, panel_w, panel_h, C.night_2, border, radius=True)
        rect(s, left, panel_t, Inches(0.08), panel_h, accent)
        oval(s, left + Inches(0.25), panel_t + Inches(0.25), Inches(0.46), Inches(0.46), accent)
        text(s, left + Inches(0.25), panel_t + Inches(0.34), Inches(0.46), Inches(0.22), num, size=11, bold=True, color=C.white, align=PP_ALIGN.CENTER)
        text(s, left + Inches(0.9), panel_t + Inches(0.22), Inches(4.55), Inches(0.32), head, size=15, bold=True, color=C.white)
        text(s, left + Inches(0.9), panel_t + Inches(0.62), Inches(4.55), Inches(0.45), body, size=9, color=muted)

        with Image.open(path) as im:
            iw, ih = im.size
        box_l = left + Inches(0.25)
        box_t = panel_t + Inches(1.25)
        box_w = panel_w - Inches(0.5)  # 5.3"
        # Width-fill (manual tweak): keep aspect, pin to content width.
        fw = box_w
        fh = int(float(box_w) * ih / iw)
        max_h = panel_h - Inches(2.05)
        if fh > max_h:
            fh = max_h
            fw = int(float(max_h) * iw / ih)
        add_media(s, path, box_l, box_t, fw, fh, caption=caption, accent=accent, theme="dark")

        callout_t = Inches(5.90)
        rect(s, box_l, callout_t, box_w, Inches(0.55), C.night, border, radius=True)
        text(s, box_l + Inches(0.15), callout_t + Inches(0.15), box_w - Inches(0.3), Inches(0.25), callout, size=9, bold=True, color=accent, align=PP_ALIGN.CENTER)

    footer_e2e(s, n, dark=True)


def slide_cover(prs, n: int):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, C.e8_black)
    rect(s, Inches(0.65), Inches(1.72), Inches(2.6), Pt(3), C.e8_coral)
    oval(s, Inches(11.95), Inches(0.5), Inches(0.75), Inches(0.75), C.e8_teal)
    rect(s, Inches(0), Inches(6.35), Inches(4.2), Inches(1.15), C.e8_teal)
    rect(s, Inches(4.2), Inches(6.35), W - Inches(4.2), Inches(1.15), C.e8_coral)

    add_logo(s, LOGO, Inches(0.65), Inches(0.5), Inches(2.4), Inches(1.1))
    text(s, Inches(0.65), Inches(2.05), Inches(4.6), Inches(2.2), "E2E пайплайн\nкоров", size=40, bold=True, color=C.white)
    text(s, Inches(0.65), Inches(4.35), Inches(4.5), Inches(0.9), "Форма → съёмка → загрузка → datapipe → визуализация → протокол", size=13, color=C.e8_coral_muted)
    text(s, Inches(0.65), Inches(5.15), Inches(4.2), Inches(0.4), "Korovas · data-collector · datapipe", size=11, color=C.e8_teal_muted)

    add_media(s, E2E_IMG / "form_create" / "config-steps-admin.gif", Inches(5.2), Inches(1.35), Inches(2.15), Inches(4.15), caption="Форма", accent=C.e8_coral, theme="dark", frame=False)
    add_media(s, E2E_IMG / "titul" / "view_on_form.png", Inches(7.55), Inches(1.45), Inches(2.05), Inches(4.05), caption="Съёмка", accent=C.e8_teal, theme="dark")
    add_media(s, E2E_IMG / "titul" / "visualization.png", Inches(9.85), Inches(1.55), Inches(2.95), Inches(3.85), caption="Визуализация", accent=C.e8_teal, theme="dark")

    text(s, Inches(0.75), Inches(6.55), Inches(12.0), Inches(0.45), "Форма · Съёмка · Загрузка · Datapipe · Визуализация · Протокол", size=18, bold=True, color=C.e8_black, align=PP_ALIGN.CENTER)
    footer_e2e(s, n)


def slide_overview(prs, n: int):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    title(s, "Полный цикл одного пакета", "Шесть этапов — от настройки формы до отправки в АИС «Монолит»")

    stages = [
        ("00", "Форма", "сценарий сбора\nв Git", C.coral),
        ("01", "Съёмка", "мобилка\nнесколько ракурсов", C.cobalt),
        ("02", "Загрузка", "отправка\nна сервер", C.coral),
        ("03", "Datapipe", "разметка · инференс\n· обучение", C.mint),
        ("04", "Визуализация", "keypoints\nи метрики", C.lime),
        ("05", "Протокол", "формирование\nдля АИС Монолит", C.coral),
    ]
    x = Inches(0.55)
    y = Inches(2.15)
    w = Inches(1.95)
    for i, (num, head, body, color) in enumerate(stages):
        rect(s, x, y, w, Inches(1.55), C.white, C.line, radius=True)
        pill(s, x + Inches(0.15), y + Inches(0.18), Inches(0.62), num, color, C.ink if color == C.lime else C.white)
        text(s, x + Inches(0.15), y + Inches(0.62), w - Inches(0.3), Inches(0.32), head, size=14, bold=True)
        text(s, x + Inches(0.15), y + Inches(0.98), w - Inches(0.3), Inches(0.55), body, size=8, color=C.muted)
        if i < len(stages) - 1:
            arrow(s, x + w, y + Inches(0.78), x + w + Inches(0.12), y + Inches(0.78), C.ink)
        x += w + Inches(0.14)

    cards = [
        ("data-collector", "Форма в Git, сбор пакета офлайн, загрузка и визуализация в админке", C.cobalt),
        ("datapipe", "Пакеты → CVAT → инференс → обучение → боевая модель", C.mint),
        ("Korovas", "21 ключевая точка, измерения, протокол для АИС «Монолит»", C.coral),
    ]
    x = Inches(0.65)
    y = Inches(4.35)
    for head, body, color in cards:
        card(s, x, y, Inches(3.85), Inches(1.15), head, body, color)
        x += Inches(4.05)
    footer_e2e(s, n)


def slide_form_create(prs, n: int):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, C.night)
    title(s, "0. Создание формы", "Первый шаг: настраиваем сценарий сбора в админке и Git", dark=True)

    from pptx.dml.color import RGBColor
    muted = RGBColor(0xAE, 0xB7, 0xC7)

    steps = [
        ("1", "Проект и Git", "Создаём проект, deploy key, привязка репозитория", C.coral),
        ("2", "Сценарий сбора", "collector/config.json — шаги, поля, ракурсы", C.cobalt),
        ("3", "Визуализация", "collector/viz.json — слои в админке", C.mint),
        ("4", "Пайплайн", "collector/pipeline.json — webhook datapipe", C.lime),
    ]
    y = Inches(2.0)
    for num, head, body, color in steps:
        rect(s, Inches(0.75), y, Inches(5.15), Inches(0.95), C.night_2, RGBColor(0x2D, 0x36, 0x48), radius=True)
        rect(s, Inches(0.75), y, Inches(0.08), Inches(0.95), color)
        oval(s, Inches(0.95), y + Inches(0.22), Inches(0.42), Inches(0.42), color)
        text(s, Inches(0.95), y + Inches(0.3), Inches(0.42), Inches(0.25), num, size=12, bold=True, color=C.white, align=PP_ALIGN.CENTER)
        text(s, Inches(1.55), y + Inches(0.14), Inches(4.1), Inches(0.3), head, size=13, bold=True, color=C.white)
        text(s, Inches(1.55), y + Inches(0.48), Inches(4.1), Inches(0.38), body, size=9, color=muted)
        y += Inches(1.08)

    text(s, Inches(0.75), Inches(6.2), Inches(5.15), Inches(0.45), "После push в Git оператор видит обновлённый сценарий в мобилке", size=10, color=muted)
    add_media(s, E2E_IMG / "form_create" / "config-steps-admin.gif", Inches(6.1), Inches(1.75), Inches(4.85), Inches(4.55), caption="Редактор формы", accent=C.coral, theme="dark", frame=False)
    add_media(s, E2E_IMG / "titul" / "view_on_form.png", Inches(11.15), Inches(2.0), Inches(1.75), Inches(4.15), caption="В приложении", accent=C.mint, theme="dark")
    footer_e2e(s, n, dark=True)


def slide_shooting(prs, n: int):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, C.night)
    title(s, "1. Съёмка пакета", "Оператор проходит сценарий — пакет сохраняется на устройстве", dark=True)

    from pptx.dml.color import RGBColor
    muted = RGBColor(0xAE, 0xB7, 0xC7)

    flow = [
        ("1", "Поля и ракурсы", "ID коровы, параметры, инструкции по съёмке", C.cobalt),
        ("2", "Фото + проверка", "Несколько ракурсов, контроль качества кадра", C.coral),
        ("3", "Сохранение", "Пакет готов к загрузке на сервер", C.mint),
    ]
    y = Inches(2.05)
    for num, head, body, color in flow:
        rect(s, Inches(0.75), y, Inches(4.85), Inches(1.05), C.night_2, RGBColor(0x2D, 0x36, 0x48), radius=True)
        rect(s, Inches(0.75), y, Inches(0.08), Inches(1.05), color)
        pill(s, Inches(0.95), y + Inches(0.18), Inches(0.5), num, color, C.white)
        text(s, Inches(1.6), y + Inches(0.16), Inches(3.8), Inches(0.3), head, size=13, bold=True, color=C.white)
        text(s, Inches(1.6), y + Inches(0.5), Inches(3.8), Inches(0.42), body, size=9, color=muted)
        y += Inches(1.18)

    add_media(s, E2E_IMG / "mobile-app" / "flutter-steps-app.gif", Inches(5.95), Inches(1.85), Inches(4.15), Inches(4.85), caption="Сценарий в приложении", accent=C.cobalt, theme="dark", frame=False)
    add_media(s, E2E_IMG / "titul" / "view_on_form.png", Inches(10.35), Inches(2.15), Inches(2.45), Inches(4.55), caption="Ракурсы съёмки", accent=C.mint, theme="dark")
    footer_e2e(s, n, dark=True)


def slide_upload(prs, n: int):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, C.night)
    title(s, "2. Загрузка на сервер", "Пошаговая отправка пакета на сервер", dark=True)

    from pptx.dml.color import RGBColor
    muted = RGBColor(0xAE, 0xB7, 0xC7)
    divider = RGBColor(0x2B, 0x34, 0x45)

    steps = [
        ("1", "Сессия", "создаём сессию загрузки для пакета"),
        ("2", "Файлы", "отправляем каждое фото отдельным запросом"),
        ("3", "Манифест", "передаём payload.json и метаданные"),
        ("4", "Commit", "пакет принят — запускается webhook datapipe"),
    ]
    y = Inches(2.0)
    for i, (num, head, body) in enumerate(steps):
        color = [C.cobalt, C.coral, C.mint, C.lime][i]
        pill(s, Inches(0.8), y, Inches(0.55), num, color, C.ink if color == C.lime else C.white)
        text(s, Inches(1.55), y + Inches(0.02), Inches(3.5), Inches(0.3), head, size=14, bold=True, color=C.white)
        text(s, Inches(1.55), y + Inches(0.34), Inches(3.8), Inches(0.3), body, size=10, color=muted)
        rect(s, Inches(0.8), y + Inches(0.72), Inches(4.55), Pt(1), divider)
        y += Inches(0.95)

    text(s, Inches(0.8), Inches(5.85), Inches(4.55), Inches(0.5), "Оператор инициирует отправку с вкладки «Сервер»", size=10, color=muted)

    add_media(s, E2E_IMG / "upload-to-server" / "mobile-upload-screen.png", Inches(5.65), Inches(1.85), Inches(2.85), Inches(4.75), caption="Вкладка «Сервер»", accent=C.coral, theme="dark")
    add_media(s, E2E_IMG / "upload-to-server" / "admin-packages.gif", Inches(8.75), Inches(1.95), Inches(3.95), Inches(4.55), caption="Список пакетов", accent=C.lime, theme="dark", frame=False)
    footer_e2e(s, n, dark=True)


def slide_datapipe_flow(prs, n: int):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, C.night)
    title(s, "3. Пайплайн datapipe", "Два контура: обработка каждого пакета и цикл обновления модели", dark=True)

    from pptx.dml.color import RGBColor
    muted = RGBColor(0xAE, 0xB7, 0xC7)
    border = RGBColor(0x2D, 0x36, 0x48)

    rect(s, Inches(0.75), Inches(1.95), Inches(11.85), Inches(0.62), C.night_2, border, radius=True)
    text(s, Inches(1.0), Inches(2.1), Inches(11.35), Inches(0.35), "Commit пакета  →  webhook  →  стадия 0 запускается автоматически", size=13, bold=True, color=C.white, align=PP_ALIGN.CENTER)

    stages = [
        ("0", "Пакеты", "На каждый новый пакет:\nфото из S3 → CVAT →\nпервичный инференс в БД", C.cobalt),
        ("1", "Аннотация", "Периодически:\nразметка из CVAT →\ndataset train/val/test", C.mint),
        ("2", "Обучение", "По накопленным данным:\ntrain keypoints-модели,\nвыбор лучшей по метрикам", C.lime),
        ("3", "Prod", "Боевая модель на пакетах:\nинференс → БД →\nслой в визуализации", C.coral),
    ]
    x = Inches(0.75)
    y = Inches(2.95)
    w = Inches(2.85)
    for i, (num, head, body, color) in enumerate(stages):
        rect(s, x, y, w, Inches(2.55), C.night_2, border, radius=True)
        rect(s, x, y, Inches(0.08), Inches(2.55), color)
        oval(s, x + Inches(0.22), y + Inches(0.28), Inches(0.52), Inches(0.52), color)
        text(s, x + Inches(0.22), y + Inches(0.36), Inches(0.52), Inches(0.28), num, size=14, bold=True, color=C.white, align=PP_ALIGN.CENTER)
        text(s, x + Inches(0.22), y + Inches(0.95), w - Inches(0.4), Inches(0.38), head, size=18, bold=True, color=C.white)
        text(s, x + Inches(0.22), y + Inches(1.42), w - Inches(0.4), Inches(0.95), body, size=9, color=muted)
        if i < len(stages) - 1:
            arrow(s, x + w, y + Inches(1.28), x + w + Inches(0.18), y + Inches(1.28), C.e8_teal_muted)
        x += w + Inches(0.2)

    text(s, Inches(0.75), Inches(5.75), Inches(11.85), Inches(0.55), "Стадия 0 — на каждый пакет после загрузки  ·  Стадии 1–3 — цикл ML-инженера для обновления модели", size=10, color=muted, align=PP_ALIGN.CENTER)
    footer_e2e(s, n, dark=True)


def slide_datapipe_stage(prs, n: int, *, stage_num: str, headline: str, subtitle: str, bullets: list[str], folder: str, accent, dark: bool = False):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, C.night if dark else C.paper)
    title(s, headline, subtitle, dark=dark)
    pill(
        s,
        Inches(11.25),
        Inches(0.55),
        Inches(1.25),
        f"ЭТАП {stage_num}",
        accent,
        C.ink if accent == C.lime else C.white,
    )

    from pptx.dml.color import RGBColor
    bullet_color = C.muted if not dark else RGBColor(0xAE, 0xB7, 0xC7)
    y = Inches(2.0)
    for item in bullets:
        rect(s, Inches(0.75), y + Inches(0.08), Inches(0.1), Inches(0.1), accent)
        text(s, Inches(1.0), y, Inches(5.2), Inches(0.55), item, size=10, color=bullet_color)
        y += Inches(0.54)

    add_media(
        s,
        _resolve_media(DATAPIPE_IMG / folder),
        Inches(6.35),
        Inches(1.9),
        Inches(6.35),
        Inches(4.75),
        caption=f"Стадия {stage_num}",
        accent=accent,
        theme="dark" if dark else "light",
        placeholder=f"Скриншот стадии {stage_num} — добавить",
    )
    footer_e2e(s, n, dark=dark)


def slide_visualization(prs, n: int):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, C.night)
    title(s, "4. Визуализация в админке", "Отображение keypoints и метрик поверх фото пакета", dark=True)

    from pptx.dml.color import RGBColor
    muted = RGBColor(0xAE, 0xB7, 0xC7)
    divider = RGBColor(0x2B, 0x34, 0x45)

    layers = [
        ("GT эксперт", "Размеры, описания, признаки — заполняет эксперт", C.coral),
        ("GT разметка", "Keypoints из CVAT — 21 точка", C.mint),
        ("Инференс", "Результат модели: keypoints и метрики", C.lime),
        ("CVAT", "Ссылка на задачу разметки", C.cobalt),
        ("Экспорт", "Скачать фото с наложенной разметкой", C.coral),
    ]
    y = Inches(1.95)
    for tag, body, color in layers:
        pill(s, Inches(0.75), y, Inches(1.3), tag, color, C.ink if color == C.lime else C.white)
        text(s, Inches(2.2), y + Inches(0.04), Inches(2.3), Inches(0.42), body, size=9, color=C.white)
        rect(s, Inches(0.75), y + Inches(0.46), Inches(3.75), Pt(1), divider)
        y += Inches(0.56)

    text(s, Inches(0.75), Inches(4.95), Inches(3.75), Inches(0.4), "GT эксперт ≠ GT разметка — разные источники данных", size=8, color=muted)

    rx = Inches(4.65)
    rw = Inches(8.0)
    add_media(s, E2E_IMG / "visualization-admin-panel" / "visual.png", rx, Inches(1.7), rw, Inches(2.85), caption="Просмотр пакета", accent=C.lime, theme="dark")

    examples = [
        ("gt-expert.png", "GT эксперт", C.coral, rx),
        ("gt.png", "GT разметка", C.mint, rx + Inches(2.65)),
        ("inference.png", "Инференс", C.lime, rx + Inches(5.3)),
    ]
    ew = Inches(2.5)
    eh = Inches(1.9)
    ey = Inches(4.75)
    for fname, cap, accent, x in examples:
        add_media(s, E2E_IMG / "visualization-admin-panel" / fname, x, ey, ew, eh, caption=cap, accent=accent, theme="dark")
    footer_e2e(s, n, dark=True)


def slide_protocol(prs, n: int):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    title(s, "5. Формирование протокола", "Подготовка данных для отправки в АИС «Монолит»")

    flow = [
        ("1", "Измерения", "21 keypoint → 15 параметров → 9 признаков", C.cobalt),
        ("2", "Проверка", "Сверка эксперта, модели и VLM", C.mint),
        ("3", "Протокол", "Формирование итогового документа бонитировки", C.coral),
        ("4", "Монолит", "Отправка протокола в АИС «Монолит»", C.lime),
    ]
    x = Inches(0.75)
    y = Inches(2.05)
    for num, head, body, color in flow:
        rect(s, x, y, Inches(5.35), Inches(0.88), C.white, C.line, radius=True)
        rect(s, x, y, Inches(0.08), Inches(0.88), color)
        oval(s, Inches(0.95), y + Inches(0.2), Inches(0.38), Inches(0.38), color)
        text(s, Inches(0.95), y + Inches(0.27), Inches(0.38), Inches(0.22), num, size=11, bold=True, color=C.white, align=PP_ALIGN.CENTER)
        text(s, Inches(1.5), y + Inches(0.12), Inches(4.4), Inches(0.28), head, size=13, bold=True)
        text(s, Inches(1.5), y + Inches(0.42), Inches(4.4), Inches(0.35), body, size=9, color=C.muted)
        y += Inches(0.98)

    text(s, Inches(0.75), Inches(6.05), Inches(5.35), Inches(0.45), "Качество съёмки и разметки напрямую влияет на точность протокола", size=10, color=C.muted)

    add_media(s, E2E_IMG / "visualization-admin-panel" / "example3.png", Inches(6.55), Inches(1.85), Inches(6.15), Inches(4.75), caption="Метрики на пакете", accent=C.coral)
    footer_e2e(s, n)


def slide_flow_recap(prs, n: int):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    rect(s, Inches(0), Inches(0), W, Inches(1.15), C.night)
    text(s, Inches(0.75), Inches(0.32), Inches(10.0), Inches(0.45), "Сквозной сценарий на стенде", size=26, bold=True, color=C.white)
    text(s, Inches(0.75), Inches(0.78), Inches(11.5), Inches(0.3), "Чек-лист для прохождения полного цикла одного пакета", size=11, color=C.muted)

    checklist = [
        ("00", "Создать форму", "config.json, viz.json, pipeline.json в Git", C.coral),
        ("01", "Снять пакет", "пройти сценарий в мобильном приложении", C.cobalt),
        ("02", "Загрузить на сервер", "сессия → файлы → манифест → commit", C.mint),
        ("03", "Прогнать datapipe", "пакеты → аннотация → обучение → prod", C.lime),
        ("04", "Открыть визуализацию", "keypoints и метрики в админке", C.cobalt),
        ("05", "Сформировать протокол", "отправка в АИС «Монолит»", C.coral),
    ]
    y = Inches(1.55)
    for num, head, body, color in checklist:
        rect(s, Inches(0.75), y, Inches(11.6), Inches(0.68), C.white, C.line, radius=True)
        rect(s, Inches(0.75), y, Inches(0.08), Inches(0.68), color)
        oval(s, Inches(0.95), y + Inches(0.14), Inches(0.38), Inches(0.38), color)
        text(s, Inches(0.95), y + Inches(0.2), Inches(0.38), Inches(0.2), num, size=9, bold=True, color=C.white, align=PP_ALIGN.CENTER)
        text(s, Inches(1.5), y + Inches(0.12), Inches(3.35), Inches(0.45), head, size=14, bold=True)
        text(s, Inches(5.0), y + Inches(0.14), Inches(7.1), Inches(0.42), body, size=10, color=C.muted)
        y += Inches(0.74)

    footer_e2e(s, n)


def build() -> Path:
    prs = new_deck()
    n = 1

    slide_cover(prs, n); n += 1
    slide_overview(prs, n); n += 1
    slide_form_create(prs, n); n += 1
    slide_video_desktop(
        prs, n,
        "Видео-каст: создание формы",
        "Заполнение полей сценария в админке",
        VIDEO / "form_create_fill_field.mp4",
        POSTERS / "form_create_fill_field.jpg",
        vw=1920, vh=1080,
        layout=LAYOUT_VIDEO_FORM,
    ); n += 1
    slide_users(prs, n); n += 1
    slide_shooting(prs, n); n += 1
    slide_upload(prs, n); n += 1
    slide_video_phone(
        prs, n,
        "Видео-каст: создание пакета",
        "Съёмка пакета в мобильном приложении",
        VIDEO / "pkg_mobile.mp4",
        POSTERS / "pkg_mobile.jpg",
    ); n += 1
    slide_datapipe_flow(prs, n); n += 1

    stages = [
        ("0", "3.1 Стадия 0: пакеты", "Пакет из админ-БД → разметка и первичный инференс", [
            "Синхронизация пакетов и файлов из БД проекта",
            "Скачивание фото из хранилища S3",
            "Первичный инференс через Gradio — результат в БД",
            "Создание задач в CVAT, привязка фото пакета",
            "Запись ссылок на CVAT и аннотаций обратно в БД",
        ], "stage-0-packages", C.cobalt),
        ("1", "3.2 Стадия 1: аннотация", "Выгрузка разметки из CVAT → подготовка обучающей выборки", [
            "Сбор аннотаций из всех CVAT-проектов",
            "Скачивание изображений из S3 или CVAT",
            "Детектор рамки животного + резервный SAM по точкам",
            "Связка keypoints с рамкой → эталонная разметка",
            "Разбиение на train/val/test",
        ], "stage-1-annotation", C.mint),
        ("2", "3.3 Стадия 2: обучение", "Обучение модели keypoints", [
            "Заморозка датасета для воспроизводимого обучения",
            "Обучение модели с аугментациями",
            "Проверка на отложенной выборке, подсчёт метрик",
            "Выбор лучшей модели по val-метрикам",
        ], "stage-2-train", C.lime, True),
        ("3", "3.4 Стадия 3: боевая модель", "Инференс prod-модели и публикация результатов", [
            "Загрузка весов боевой модели",
            "Инференс на пакетах — результат в БД проекта",
            "Анализ ошибок по предсказаниям",
            "Слой «Инференс» виден в визуализации админки",
        ], "stage-4-prod", C.coral, True),
    ]

    for stage_num, headline, subtitle, bullets, folder, accent, *rest in stages:
        slide_datapipe_stage(prs, n, stage_num=stage_num, headline=headline, subtitle=subtitle, bullets=bullets, folder=folder, accent=accent, dark=bool(rest and rest[0]))
        n += 1

    slide_video_placeholder(prs, n, "Видео-каст: пайплайн datapipe", "Полный проход пакета через трубу"); n += 1
    slide_visualization(prs, n); n += 1
    slide_video_desktop(
        prs, n,
        "Видео-каст: просмотр пакета",
        "Визуализация пакета в админке — слои разметки и метрики",
        VIDEO / "visualization_pkg.mp4",
        POSTERS / "visualization_pkg.jpg",
        vw=1918, vh=946,
        layout=LAYOUT_VIDEO_VIZ,
    ); n += 1
    slide_protocol(prs, n); n += 1
    slide_flow_recap(prs, n)

    target = OUT
    try:
        prs.save(target)
    except PermissionError:
        target = OUT.with_name(f"{OUT.stem}-generated{OUT.suffix}")
        prs.save(target)
        print(f"Note: {OUT.name} is open — saved as {target.name}")
    return target


if __name__ == "__main__":
    print(f"Saved: {build()}")
