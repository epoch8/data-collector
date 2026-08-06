"""
Generate Architecture presentation — diagram-first, based on specs + datapipe.

Run: python docs/architecture/generate_presentation.py
"""

from __future__ import annotations

import importlib.util
import os
import sys
from pathlib import Path

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent
REPO = ROOT.parent.parent
BUSINESS = ROOT.parent / "business"
LOGO = REPO / "e8-team-logo-1024.png"
OUT = ROOT / "Architecture.pptx"
VIDEO = ROOT / "video"
POSTERS = VIDEO / "posters"

spec = importlib.util.spec_from_file_location("gen_training", BUSINESS / "generate_training_presentation.py")
gen = importlib.util.module_from_spec(spec)
sys.modules["gen_training"] = gen
spec.loader.exec_module(gen)

C = gen.C
W = gen.W
TOTAL = 13

bg = gen.bg
text = gen.text
rect = gen.rect
oval = gen.oval
pill = gen.pill
add_logo = gen.add_logo
new_deck = gen.new_deck

BG = RGBColor(0x09, 0x0D, 0x16)
PANEL = RGBColor(0x0F, 0x14, 0x19)
CARD = RGBColor(0x1E, 0x2A, 0x3D)
CARD2 = RGBColor(0x13, 0x1A, 0x2A)
BLUE = RGBColor(0x5B, 0x8D, 0xEF)
GREEN = RGBColor(0x3D, 0xCC, 0x85)
GOLD = RGBColor(0xD4, 0xA8, 0x4B)
PURPLE = RGBColor(0x9B, 0x7E, 0xD0)
CORAL = RGBColor(0xE8, 0x7A, 0x7E)
TEAL = RGBColor(0x2D, 0xD4, 0xBF)
FIRE = RGBColor(0xF9, 0xAB, 0x00)
TEXT = RGBColor(0xE8, 0xEA, 0xED)
MUTED = RGBColor(0x9C, 0xA3, 0xAF)
BORDER = RGBColor(0x3A, 0x45, 0x58)
WHITE = C.white


def footer(slide, n: int) -> None:
    text(slide, Inches(0.55), Inches(7.15), Inches(5.5), Inches(0.22), "Epoch8 · Architecture", size=8, color=MUTED)
    text(slide, Inches(12.0), Inches(7.15), Inches(0.8), Inches(0.22), f"{n:02d}/{TOTAL:02d}", size=8, color=MUTED, align=PP_ALIGN.RIGHT)


def headline(slide, title_text: str, subtitle: str | None = None) -> None:
    text(slide, Inches(0.55), Inches(0.28), Inches(12.2), Inches(0.45), title_text, size=24, bold=True, color=TEXT)
    if subtitle:
        text(slide, Inches(0.55), Inches(0.72), Inches(12.2), Inches(0.28), subtitle, size=11, color=MUTED)


def box(slide, l, t, w, h, fill, line, title, sub="", *, ts=12, ss=9, align=PP_ALIGN.CENTER):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.color.rgb = line
    s.line.width = Pt(1.5)
    tf = s.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Pt(8)
    tf.margin_right = Pt(8)
    tf.margin_top = Pt(4)
    tf.margin_bottom = Pt(4)
    p0 = tf.paragraphs[0]
    p0.text = title
    p0.font.size = Pt(ts)
    p0.font.bold = True
    p0.font.color.rgb = TEXT
    p0.alignment = align
    if sub:
        p1 = tf.add_paragraph()
        p1.text = sub
        p1.font.size = Pt(ss)
        p1.font.color.rgb = MUTED
        p1.alignment = align
        p1.space_before = Pt(3)
    return s


def lane(slide, l, t, w, h, stroke, label=""):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = PANEL
    s.line.color.rgb = stroke
    s.line.width = Pt(2)
    if label:
        text(slide, l + Inches(0.15), t + Inches(0.08), w - Inches(0.3), Inches(0.28), label, size=11, bold=True, color=stroke)
    return s


def h_arrow(slide, x1, y, x2, color):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y, x2, y)
    c.line.color.rgb = color
    c.line.width = Pt(2)
    return c


def v_arrow(slide, x, y1, y2, color):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x, y1, x, y2)
    c.line.color.rgb = color
    c.line.width = Pt(2)
    return c


def label(slide, l, t, w, value, *, size=9, color=MUTED, align=PP_ALIGN.CENTER):
    text(slide, l, t, w, Inches(0.25), value, size=size, color=color, align=align)


def ensure_video_poster(movie: Path, poster: Path, *, at_sec: float = 2.0) -> Path | None:
    """Кадр для превью в PowerPoint (без этого python-pptx ставит иконку «колонки»)."""
    if poster.exists() and poster.stat().st_size > 10_000:
        return poster
    if not movie.exists():
        return None
    try:
        import io
        import subprocess

        import imageio_ffmpeg
        from PIL import Image

        ffmpeg = imageio_ffmpeg.get_ffmpeg_exe()
        cmd = [
            ffmpeg,
            "-y",
            "-ss",
            str(at_sec),
            "-i",
            str(movie),
            "-frames:v",
            "1",
            "-f",
            "image2pipe",
            "-vcodec",
            "png",
            "-",
        ]
        proc = subprocess.run(cmd, capture_output=True, check=False)
        if proc.returncode != 0 or not proc.stdout:
            return None
        poster.parent.mkdir(parents=True, exist_ok=True)
        img = Image.open(io.BytesIO(proc.stdout))
        # чуть уменьшаем превью для pptx, исходный mp4 не трогаем
        max_w = 1920
        if img.width > max_w:
            ratio = max_w / img.width
            img = img.resize((max_w, int(img.height * ratio)), Image.Resampling.LANCZOS)
        img.save(poster, format="PNG", optimize=True)
        return poster
    except Exception as e:  # noqa: BLE001
        print(f"Warning: poster for {movie.name}: {e}")
        return None


def store_videos_uncompressed(pptx_path: Path) -> None:
    """python-pptx кладёт mp4 с ZIP_DEFLATED — перепаковываем media без сжатия."""
    import shutil
    import tempfile
    import zipfile

    fd, tmp_name = tempfile.mkstemp(suffix=".pptx")
    os.close(fd)
    tmp = Path(tmp_name)
    try:
        with zipfile.ZipFile(pptx_path, "r") as zin, zipfile.ZipFile(tmp, "w") as zout:
            for info in zin.infolist():
                data = zin.read(info.filename)
                out = zipfile.ZipInfo(filename=info.filename, date_time=info.date_time)
                out.external_attr = info.external_attr
                out.create_system = info.create_system
                lower = info.filename.lower()
                if lower.endswith((".mp4", ".mov", ".webm", ".avi", ".m4v")):
                    out.compress_type = zipfile.ZIP_STORED
                else:
                    out.compress_type = zipfile.ZIP_DEFLATED
                zout.writestr(out, data)
        shutil.move(str(tmp), str(pptx_path))
    finally:
        if tmp.exists():
            tmp.unlink(missing_ok=True)


def probe_video_size(movie: Path) -> tuple[int, int]:
    """Return (width, height) of video; fallback 16:9."""
    import re
    import subprocess

    try:
        import imageio_ffmpeg

        ffmpeg = imageio_ffmpeg.get_ffmpeg_exe()
        proc = subprocess.run(
            [ffmpeg, "-i", str(movie), "-f", "null", "-"],
            capture_output=True,
            check=False,
        )
        err = (proc.stderr or b"").decode("utf-8", "replace")
        m = re.search(r"Video:.*?,\s*(\d{2,5})x(\d{2,5})", err)
        if m:
            return int(m.group(1)), int(m.group(2))
    except Exception:  # noqa: BLE001
        pass
    try:
        from PIL import Image

        poster = POSTERS / f"{movie.stem}.png"
        if poster.exists():
            with Image.open(poster) as im:
                return im.size
    except Exception:  # noqa: BLE001
        pass
    return 1920, 1080


def fit_box(src_w: int, src_h: int, max_w, max_h) -> tuple:
    """Fit inside max box, keep aspect ratio. Returns (w, h) in same units as max_*."""
    scale = min(float(max_w) / src_w, float(max_h) / src_h)
    return int(src_w * scale), int(src_h * scale)


def enable_video_fullscreen(slide) -> None:
    """В слайд-шоу Play открывает видео на весь экран (OOXML p:video/@fullScrn)."""
    for video in slide._element.xpath(".//p:video"):
        video.set("fullScrn", "1")


def slide_video_scenario(
    prs,
    n: int,
    title_text: str,
    subtitle: str,
    steps: list[str],
    *,
    note: str | None = None,
    video_hint: str = "▶  Вставить видео",
    movie: Path | None = None,
    poster: Path | None = None,
) -> None:
    """Слайд-сценарий: шаги слева, видео справа (или плейсхолдер, если файла нет)."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, BG)
    headline(s, title_text, subtitle)

    # left: scenario steps
    left_w = Inches(5.6)
    rect(s, Inches(0.45), Inches(1.15), left_w, Inches(5.7), CARD2, BORDER, radius=True)
    text(s, Inches(0.65), Inches(1.3), Inches(5.2), Inches(0.32), "Сценарий", size=12, bold=True, color=TEAL)
    y = Inches(1.7)
    for i, step in enumerate(steps, 1):
        pill(s, Inches(0.7), y, Inches(0.38), str(i), BLUE, WHITE)
        text(s, Inches(1.2), y + Inches(0.02), Inches(4.55), Inches(0.42), step, size=11, color=TEXT)
        y += Inches(0.42)
        if y > Inches(6.4):
            break
    if note:
        text(s, Inches(0.65), Inches(6.35), Inches(5.2), Inches(0.4), note, size=9, color=MUTED)

    # right: video area — fit without stretch
    area_l, area_t = Inches(6.3), Inches(1.15)
    area_w, area_h = Inches(6.55), Inches(5.7)
    if movie is not None and movie.exists():
        poster_path = poster if poster is not None else POSTERS / f"{movie.stem}.png"
        frame = ensure_video_poster(movie, poster_path)
        vw_px, vh_px = probe_video_size(movie)
        vid_w, vid_h = fit_box(vw_px, vh_px, area_w, area_h)
        vid_l = area_l + (area_w - vid_w) // 2
        vid_t = area_t + (area_h - vid_h) // 2
        s.shapes.add_movie(
            str(movie),
            vid_l,
            vid_t,
            vid_w,
            vid_h,
            poster_frame_image=str(frame) if frame else None,
            mime_type="video/mp4",
        )
        enable_video_fullscreen(s)
        text(
            s,
            area_l,
            area_t + area_h - Inches(0.28),
            area_w,
            Inches(0.25),
            "Слайд-шоу → ▶  ·  полный экран",
            size=9,
            color=MUTED,
            align=PP_ALIGN.CENTER,
        )
    else:
        rect(s, area_l, area_t, area_w, area_h, CARD, BORDER, radius=True)
        text(
            s,
            area_l,
            area_t + area_h / 2 - Inches(0.35),
            area_w,
            Inches(0.4),
            video_hint,
            size=18,
            bold=True,
            color=TEXT,
            align=PP_ALIGN.CENTER,
        )
        text(
            s,
            area_l + Inches(0.4),
            area_t + area_h / 2 + Inches(0.15),
            area_w - Inches(0.8),
            Inches(0.5),
            "В PowerPoint: Вставка → Видео → Этот компьютер",
            size=10,
            color=MUTED,
            align=PP_ALIGN.CENTER,
        )
    footer(s, n)


# ─── slides ───────────────────────────────────────────────────────────────────


def slide_cover(prs, n):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, C.e8_black)
    rect(s, Inches(0), Inches(6.35), Inches(4.2), Inches(1.15), C.e8_teal)
    rect(s, Inches(4.2), Inches(6.35), W - Inches(4.2), Inches(1.15), C.e8_coral)
    add_logo(s, LOGO, Inches(0.65), Inches(0.45), Inches(2.2), Inches(1.0))
    rect(s, Inches(0.65), Inches(1.65), Inches(2.4), Pt(3), C.e8_coral)
    text(s, Inches(0.65), Inches(1.95), Inches(7), Inches(1.6), "Архитектура", size=44, bold=True, color=WHITE)
    text(s, Inches(0.65), Inches(3.6), Inches(9), Inches(0.4), "Data Collector · стек · поток данных · пайплайны", size=16, color=C.e8_coral_muted)

    nodes = [("Flutter", BLUE), ("Django", GREEN), ("Хранилище", GOLD), ("Datapipe", TEAL), ("Визуализация", PURPLE)]
    x = Inches(0.65)
    y = Inches(4.5)
    bw = Inches(1.9)
    for i, (name, col) in enumerate(nodes):
        box(s, x, y, bw, Inches(0.7), CARD, col, name, ts=13)
        if i < len(nodes) - 1:
            h_arrow(s, x + bw, y + Inches(0.35), x + bw + Inches(0.25), col)
        x += bw + Inches(0.25)

    text(s, Inches(0.75), Inches(6.6), Inches(12), Inches(0.4), "схемы по specs/ и datapipe", size=14, bold=True, color=C.e8_black, align=PP_ALIGN.CENTER)
    footer(s, n)


def slide_stack(prs, n):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, BG)
    headline(s, "Стек: 1 платформа · N проектов")

    lane(s, Inches(0.45), Inches(1.1), Inches(12.4), Inches(3.2), BLUE, "ПЛАТФОРМА  —  data-collector")

    label(s, Inches(0.65), Inches(1.45), Inches(3.4), "КЛИЕНТЫ", size=9, color=MUTED, align=PP_ALIGN.LEFT)
    box(s, Inches(0.65), Inches(1.75), Inches(1.7), Inches(1.05), RGBColor(0x2A, 0x24, 0x38), PURPLE, "Flutter", "мобилка\nAPI /v1", ts=12, ss=9)
    box(s, Inches(2.5), Inches(1.75), Inches(1.7), Inches(1.05), RGBColor(0x3D, 0x34, 0x20), GOLD, "Staff", "админ\nпроекты, юзеры", ts=12, ss=9)
    box(s, Inches(0.65), Inches(2.95), Inches(3.55), Inches(0.95), RGBColor(0x1A, 0x2E, 0x28), GREEN, "Client-admin", "пакеты и визуализация проекта", ts=12, ss=10)

    label(s, Inches(4.6), Inches(1.45), Inches(3.5), "СЕРВЕР", size=9, color=MUTED, align=PP_ALIGN.LEFT)
    box(s, Inches(4.6), Inches(1.75), Inches(4.0), Inches(2.15), CARD, BLUE, "django_server", "мобильный API  /v1\nадминка  /ui\nконфиг, медиа, плагины viz", ts=14, ss=11)

    box(s, Inches(9.0), Inches(1.75), Inches(3.5), Inches(0.95), RGBColor(0x3D, 0x30, 0x20), FIRE, "Firebase", "вход и токен", ts=13, ss=10)
    box(s, Inches(9.0), Inches(2.9), Inches(3.5), Inches(1.0), RGBColor(0x1A, 0x33, 0x29), GREEN, "Каталог Django", "проекты, юзеры, ключи Git\nпакетов здесь нет", ts=12, ss=9)

    lane(s, Inches(0.45), Inches(4.5), Inches(12.4), Inches(2.35), TEAL, "ПРОЕКТ  —  ×N  (у каждого своё)")

    items = [
        ("Git", "сценарий сбора\nconfig · viz · pipeline", BLUE),
        ("Кэш Git", "копия репо\nна сервере", BLUE),
        ("БД проекта", "PostgreSQL\nпакеты и keypoints", PURPLE),
        ("Файлы", "S3 / GCS\nфото пакетов", GREEN),
        ("… проект N", "свой Git · БД · бакет", BORDER),
    ]
    x = Inches(0.65)
    for title, sub, col in items:
        box(s, x, Inches(5.0), Inches(2.25), Inches(1.5), CARD2, col, title, sub, ts=13, ss=10)
        x += Inches(2.4)

    footer(s, n)


def slide_inside(prs, n):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, BG)
    headline(s, "Как устроено изнутри")

    lane(s, Inches(0.5), Inches(1.15), Inches(3.5), Inches(5.5), PURPLE, "МОБИЛКА")
    for i, (t, sub) in enumerate([
        ("Состояние", "Riverpod"),
        ("Локальная БД", "Drift · пакеты на устройстве"),
        ("Сеть", "Dio + токен Firebase"),
        ("Сборка пакета", "манифест + фото"),
        ("Отправка", "вкладка «Сервер»"),
    ]):
        box(s, Inches(0.75), Inches(1.6) + Inches(i * 0.9), Inches(3.0), Inches(0.78), CARD, PURPLE, t, sub, ts=12, ss=10)

    text(s, Inches(4.15), Inches(3.3), Inches(1.5), Inches(0.7), "HTTPS\n/v1", size=14, bold=True, color=BLUE, align=PP_ALIGN.CENTER)
    h_arrow(s, Inches(4.0), Inches(3.55), Inches(5.55), BLUE)

    lane(s, Inches(5.7), Inches(1.15), Inches(7.1), Inches(5.5), BLUE, "СЕРВЕР")
    box(s, Inches(6.0), Inches(1.65), Inches(6.5), Inches(1.05), CARD, BLUE, "API и админка", "приём пакетов · просмотр · viz", ts=14, ss=11)

    stores = [
        ("Каталог", "проекты и права\n(Postgres / SQLite)", GREEN),
        ("БД проекта", "PostgreSQL\nсессии · keypoints", PURPLE),
        ("Объектное\nхранилище", "фото в S3 / GCS", GOLD),
        ("Git", "конфиг сценария\nи viz.json", TEAL),
    ]
    for i, (t, sub, col) in enumerate(stores):
        col_i, row = i % 2, i // 2
        box(s, Inches(6.0) + col_i * Inches(3.3), Inches(3.0) + row * Inches(1.55), Inches(3.1), Inches(1.35), CARD2, col, t, sub, ts=13, ss=11)

    footer(s, n)


def slide_config_cycle(prs, n):
    """Clearer: what lives where and the one-way purpose of each step."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, BG)
    headline(s, "От сценария до пакета на сервере", "зачем каждый шаг")

    steps = [
        ("1", "Сценарий в Git", "что снимать:\nполя и шаги формы", "Staff пишет\nconfig.json", GOLD),
        ("2", "Мобилка\nполучает конфиг", "тот же сценарий\nна устройстве", "GET /v1/…/config", BLUE),
        ("3", "Оператор\nснимает пакет", "поля + фото\nлокально", "submit → пакет", PURPLE),
        ("4", "Загрузка\nна сервер", "пакет в БД\nпроекта + S3", "upload → commit", GREEN),
    ]
    x = Inches(0.45)
    for i, (num, title, sub, tag, col) in enumerate(steps):
        box(s, x, Inches(1.4), Inches(2.95), Inches(3.4), CARD, col, f"{num}.  {title}", sub, ts=14, ss=12)
        box(s, x + Inches(0.15), Inches(5.0), Inches(2.65), Inches(0.65), CARD2, col, tag, ts=11)
        if i < len(steps) - 1:
            h_arrow(s, x + Inches(2.95), Inches(3.1), x + Inches(3.2), col)
        x += Inches(3.25)

    box(
        s, Inches(0.45), Inches(5.9), Inches(12.4), Inches(0.85), PANEL, TEAL,
        "Один сценарий в Git  →  тот же сценарий в мобилке  →  пакет той же формы в БД проекта",
        ts=13,
    )
    footer(s, n)


def slide_upload(prs, n):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, BG)
    headline(s, "Протокол загрузки пакета")

    steps = [
        ("1", "Сессия", "открываем приём\nпакета", BLUE),
        ("2", "Файлы", "каждое фото\nотдельным запросом", TEAL),
        ("3", "Манифест", "JSON с полями\nи ссылками на фото", GREEN),
        ("4", "Commit", "пакет принят\n→ можно в datapipe", GOLD),
    ]
    x = Inches(0.55)
    for i, (num, title, sub, col) in enumerate(steps):
        box(s, x, Inches(1.35), Inches(2.7), Inches(2.3), CARD, col, f"{num}\n{title}", sub, ts=16, ss=12)
        if i < len(steps) - 1:
            h_arrow(s, x + Inches(2.7), Inches(2.5), x + Inches(3.0), col)
        x += Inches(3.15)

    label(s, Inches(0.55), Inches(3.9), Inches(6), "Фаза на сервере (таблица сессий)", size=11, color=MUTED, align=PP_ALIGN.LEFT)
    phases = [
        ("ждём файлы", BLUE),
        ("готово к commit", GREEN),
        ("принят", GOLD),
        ("ошибка", CORAL),
    ]
    x = Inches(0.55)
    for name, col in phases:
        box(s, x, Inches(4.25), Inches(2.95), Inches(0.75), CARD2, col, name, ts=13)
        x += Inches(3.15)

    label(s, Inches(0.55), Inches(5.25), Inches(6), "Статус на телефоне", size=11, color=MUTED, align=PP_ALIGN.LEFT)
    for i, (name, col) in enumerate([
        ("в очереди", MUTED),
        ("отправляется", BLUE),
        ("на сервере", GREEN),
        ("сбой", CORAL),
    ]):
        pill(s, Inches(0.55) + Inches(i * 3.1), Inches(5.6), Inches(2.8), name, col, WHITE if col != MUTED else BG)

    footer(s, n)


def slide_storage(prs, n):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, BG)
    headline(s, "Где что лежит", "четыре разных места — четыре разных смысла")

    layers = [
        ("1", "Каталог платформы", "Postgres (prod)", "Список проектов, пользователи,\nправа, ключи Git", "пакетов нет", GREEN),
        ("2", "БД проекта", "PostgreSQL", "Сессии загрузки, учёт файлов,\nkeypoints, инференс, ссылки CVAT", "данные пакетов", PURPLE),
        ("3", "Файлы пакетов", "S3 / GCS", "Сами фото: packages/…/blobs/", "байты картинок", GOLD),
        ("4", "Конфиг проекта", "Git", "Сценарий сбора, viz.json,\npipeline.json, инструкции", "что и как снимать", BLUE),
    ]
    y = Inches(1.2)
    for num, name, eng, body, tag, col in layers:
        lane(s, Inches(0.45), y, Inches(12.4), Inches(1.3), col)
        oval(s, Inches(0.65), y + Inches(0.38), Inches(0.55), Inches(0.55), col)
        text(s, Inches(0.65), y + Inches(0.48), Inches(0.55), Inches(0.35), num, size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
        text(s, Inches(1.45), y + Inches(0.2), Inches(3.2), Inches(0.4), name, size=17, bold=True, color=TEXT)
        text(s, Inches(1.45), y + Inches(0.65), Inches(3.2), Inches(0.4), eng, size=12, color=col)
        text(s, Inches(4.9), y + Inches(0.3), Inches(5.2), Inches(0.8), body, size=13, color=MUTED)
        box(s, Inches(10.3), y + Inches(0.32), Inches(2.3), Inches(0.7), CARD2, col, tag, ts=11)
        y += Inches(1.4)

    footer(s, n)


def slide_package_model(prs, n):
    """Tables in project DB — no 'on disk', no YOLO."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, BG)
    headline(s, "БД проекта: что в таблицах", "PostgreSQL · SQLAlchemy")

    lane(s, Inches(0.45), Inches(1.15), Inches(6.1), Inches(5.5), BLUE, "ПРИЁМ ПАКЕТОВ")
    upload_tables = [
        ("package_session", "сессия загрузки\nфаза: ждём файлы → принят", BLUE),
        ("uploaded_blob", "учёт каждого фото\nпуть в S3 / GCS", TEAL),
        ("package_field_change", "история правок\nманифеста в админке", GOLD),
    ]
    y = Inches(1.7)
    for name, sub, col in upload_tables:
        box(s, Inches(0.7), y, Inches(5.6), Inches(1.35), CARD, col, name, sub, ts=15, ss=12, align=PP_ALIGN.LEFT)
        y += Inches(1.5)

    lane(s, Inches(6.8), Inches(1.15), Inches(6.05), Inches(5.5), GREEN, "PIPELINE / КОРОВЫ")
    ml_tables = [
        ("cow_keypoint_annotation", "разметка keypoints (GT)\nиз CVAT", GREEN),
        ("cow_inference_result", "результат модели\nна фото пакета", CORAL),
        ("cvat_link", "ссылка на задачу\nразметки в CVAT", PURPLE),
        ("depth_map", "карта глубины\n(если есть)", TEAL),
    ]
    y = Inches(1.7)
    for name, sub, col in ml_tables:
        box(s, Inches(7.05), y, Inches(5.55), Inches(1.1), CARD2, col, name, sub, ts=13, ss=11, align=PP_ALIGN.LEFT)
        y += Inches(1.2)

    footer(s, n)


def slide_auth_roles(prs, n):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, BG)
    headline(s, "Кто куда заходит и что видит")

    # Three clear role columns
    roles = [
        (
            "Суперадмин (Staff)",
            GOLD,
            "Вход: логин Django",
            [
                "все проекты",
                "юзеры и права",
                "Git и хранилище",
                "настройка платформы",
            ],
        ),
        (
            "Админ проекта",
            GREEN,
            "Вход: Firebase → /ui/",
            [
                "только свои проекты",
                "список пакетов",
                "визуализация",
                "правки данных",
            ],
        ),
        (
            "Сборщик",
            PURPLE,
            "Вход: Firebase → мобилка",
            [
                "только выданные проекты",
                "съёмка по сценарию",
                "загрузка пакетов",
                "без чужих данных",
            ],
        ),
    ]
    x = Inches(0.45)
    for title, col, how, bullets in roles:
        lane(s, x, Inches(1.2), Inches(4.05), Inches(5.4), col)
        text(s, x + Inches(0.2), Inches(1.4), Inches(3.65), Inches(0.55), title, size=16, bold=True, color=TEXT)
        box(s, x + Inches(0.2), Inches(2.05), Inches(3.65), Inches(0.7), CARD2, col, how, ts=12)
        y = Inches(3.0)
        for b in bullets:
            rect(s, x + Inches(0.35), y + Inches(0.12), Inches(0.12), Inches(0.12), col)
            text(s, x + Inches(0.65), y, Inches(3.5), Inches(0.45), b, size=14, color=MUTED)
            y += Inches(0.7)
        x += Inches(4.25)

    footer(s, n)


def slide_local_run_video(prs, n):
    movie = VIDEO / "local run" / "create_simple_project.mp4"
    slide_video_scenario(
        prs,
        n,
        "Видео: от сценария до пакета",
        "локальный стенд · простой пример",
        [
            "Создаём репозиторий проекта в Git",
            "Создаём проект в админке и привязываем Git",
            "Поднимаем Postgres + MinIO, привязываем к проекту",
            "Делаем форму проекта",
            "Выдаём доступ к проекту",
            "Смотрим проект в мобилке",
            "Заполняем форму и отправляем пакет",
            "Смотрим пакет в админке",
            "Привязываем визуализацию",
            "Имитация инференса (запись в БД)",
            "Смотрим результат в админке",
        ],
        movie=movie,
        poster=POSTERS / "create_simple_project.png",
    )


def slide_datapipe(prs, n):
    """Based on Downloads/readme.md — stages 0..4, two data entries."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, BG)
    headline(s, "Datapipe: keypoints коров", "пакеты → CVAT → обучение → prod")

    box(
        s, Inches(0.45), Inches(1.15), Inches(12.4), Inches(0.85), CARD, TEAL,
        "Два входа:  новые пакеты из БД проекта  ·  разметка из всех проектов CVAT",
        ts=14,
    )

    stages = [
        ("0", "Пакеты", "из БД проекта:\nфото → Gradio →\nCVAT → обратно в БД", BLUE),
        ("1", "Аннотация", "все проекты CVAT:\nразметка → bbox →\ndataset train/val/test", GREEN),
        ("2", "Обучение", "заморозка датасета\ntrain keypoints\nвыбор лучшей", GOLD),
        ("3", "FiftyOne", "публикация GT\nи предсказаний", PURPLE),
        ("4", "Prod", "боевая модель\nинференс → БД\nслой в админке", CORAL),
    ]
    x = Inches(0.4)
    for i, (num, title, sub, col) in enumerate(stages):
        box(s, x, Inches(2.25), Inches(2.4), Inches(2.9), CARD, col, f"{num}  {title}", sub, ts=14, ss=11)
        if i < len(stages) - 1:
            h_arrow(s, x + Inches(2.4), Inches(3.7), x + Inches(2.55), col)
        x += Inches(2.55)

    box(s, Inches(0.45), Inches(5.45), Inches(6.1), Inches(1.25), CARD2, BLUE,
        "На каждый пакет", "стадия 0 · после commit / webhook CVAT", ts=14, ss=12)
    box(s, Inches(6.75), Inches(5.45), Inches(6.1), Inches(1.25), CARD2, CORAL,
        "Цикл модели", "стадии 1 → 2 → 3 → 4 · ML-инженер", ts=14, ss=12)

    footer(s, n)


def slide_datapipe_stage0(prs, n):
    """Stage 0 detail from readme — clearest operational path."""
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, BG)
    headline(s, "Стадия 0: пакет → CVAT → БД", "что происходит с каждым completed-пакетом")

    steps = [
        ("1", "Читаем пакеты", "package_session\nuploaded_blob", BLUE),
        ("2", "Качаем фото", "из S3\nupload-packages", TEAL),
        ("3", "Инференс", "Gradio →\ncow_inference_result", GOLD),
        ("4", "CVAT", "создаём задачу\nразметки", GREEN),
        ("5", "Пишем назад", "cvat_link +\nkeypoints GT", CORAL),
    ]
    x = Inches(0.4)
    for i, (num, title, sub, col) in enumerate(steps):
        box(s, x, Inches(1.4), Inches(2.35), Inches(2.6), CARD, col, f"{num}\n{title}", sub, ts=14, ss=11)
        if i < len(steps) - 1:
            h_arrow(s, x + Inches(2.35), Inches(2.7), x + Inches(2.55), col)
        x += Inches(2.55)

    # trigger / outputs
    box(s, Inches(0.45), Inches(4.3), Inches(6.1), Inches(2.2), CARD2, TEAL,
        "Триггер",
        "commit пакета на сервере\nили webhook CVAT → Datapipe API\n(job в acceptance / completed)",
        ts=15, ss=13, align=PP_ALIGN.LEFT)
    box(s, Inches(6.75), Inches(4.3), Inches(6.1), Inches(2.2), CARD2, GREEN,
        "В админке видно",
        "слой инференса\nссылка на CVAT\nkeypoints после разметки",
        ts=15, ss=13, align=PP_ALIGN.LEFT)

    footer(s, n)


def slide_datapipe_video(prs, n):
    movie = VIDEO / "local run" / "full_datapipe_local.mp4"
    slide_video_scenario(
        prs,
        n,
        "Видео: collector → Datapipe → CVAT",
        "локальный стенд · packages + webhook · без Gradio",
        [
            "pipeline.json: on_commit → Datapipe :8010",
            "Проверить Git в админке",
            "Commit пакета → авто stage=packages",
            "Datapipe UI: run / логи / граф",
            "CVAT: разметка cow + cow_keypoints",
            "Acceptance → Completed",
            "Webhook CVAT → снова packages",
            "Админка: слои CVAT + GT",
        ],
        note="Гайд: docs/architecture/local-run-korovas-datapipe.ru.md",
        movie=movie,
        poster=POSTERS / "full_datapipe_local.png",
    )


def slide_failure_map(prs, n):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, BG)
    headline(s, "Где искать сбой", "korovas-broken · по слою потока")

    cols = [
        ("Загрузка", BLUE, [
            "нет всех фото\n→ догрузить",
            "обрыв сети\n→ повторить",
            "нет доступа\n→ выдать проект",
        ]),
        ("Хранилище", GOLD, [
            "пакет в БД,\nфайлов нет",
            "неверный URI\nили креды S3",
            "путать каталог\nи БД проекта",
        ]),
        ("Datapipe", TEAL, [
            "не стартовал\nпосле commit",
            "не читает\nфото из S3",
            "нет задачи\nв CVAT",
        ]),
        ("Визуализация", CORAL, [
            "пустой слой\nинференса",
            "сломан\nviz.json",
            "данные в БД,\nUI пустой",
        ]),
    ]
    x = Inches(0.4)
    for title, col, items in cols:
        lane(s, x, Inches(1.2), Inches(3.1), Inches(5.4), col)
        text(s, x + Inches(0.15), Inches(1.4), Inches(2.8), Inches(0.45), title, size=16, bold=True, color=TEXT)
        y = Inches(2.05)
        for item in items:
            box(s, x + Inches(0.15), y, Inches(2.8), Inches(1.25), CARD2, col, item, ts=12)
            y += Inches(1.4)
        x += Inches(3.25)

    footer(s, n)


def build() -> Path:
    prs = new_deck()
    n = 1
    slide_cover(prs, n); n += 1
    slide_stack(prs, n); n += 1
    slide_inside(prs, n); n += 1
    slide_config_cycle(prs, n); n += 1
    slide_upload(prs, n); n += 1
    slide_storage(prs, n); n += 1
    slide_package_model(prs, n); n += 1
    slide_auth_roles(prs, n); n += 1
    slide_local_run_video(prs, n); n += 1
    slide_datapipe(prs, n); n += 1
    slide_datapipe_stage0(prs, n); n += 1
    slide_datapipe_video(prs, n); n += 1
    slide_failure_map(prs, n)

    assert n == TOTAL, f"Expected {TOTAL} slides, got {n}"

    try:
        prs.save(OUT)
    except PermissionError:
        raise SystemExit(
            f"{OUT.name} is open in PowerPoint — close it and re-run the generator."
        ) from None
    store_videos_uncompressed(OUT)
    return OUT


if __name__ == "__main__":
    print(f"Saved: {build()}")
