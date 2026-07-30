"""
Generate Datapipe intro presentation — framework overview for training.

Style aligned with docs/e2e-korovas: mix of light/dark, left copy + right screenshot,
thin accent cards — not a grid of identical cubes.

Run: python docs/datapipe/generate_presentation.py
"""

from __future__ import annotations

import importlib.util
import sys
from pathlib import Path

from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

ROOT = Path(__file__).resolve().parent
REPO = ROOT.parent.parent
BUSINESS = ROOT.parent / "business"
IMG = ROOT / "img"
PITCH = IMG / "pitch"
TAGS = IMG / "tags-demo"
LOGO = REPO / "e8-team-logo-1024.png"
OUT = ROOT / "Datapipe.pptx"

spec = importlib.util.spec_from_file_location("gen_training", BUSINESS / "generate_training_presentation.py")
gen = importlib.util.module_from_spec(spec)
sys.modules["gen_training"] = gen
spec.loader.exec_module(gen)

C = gen.C
W = gen.W
TOTAL = 11

bg = gen.bg
title = gen.title
text = gen.text
rect = gen.rect
oval = gen.oval
pill = gen.pill
card = gen.card
add_logo = gen.add_logo
new_deck = gen.new_deck

MUTED_DARK = RGBColor(0xAE, 0xB7, 0xC7)
BORDER_DARK = RGBColor(0x2D, 0x36, 0x48)
DIVIDER = RGBColor(0x2B, 0x34, 0x45)


def footer(slide, n: int, *, dark: bool = False) -> None:
    color = RGBColor(0x7C, 0x83, 0x90) if dark else C.faint
    text(slide, Inches(0.65), Inches(7.1), Inches(5.0), Inches(0.22), "Epoch8 · Datapipe", size=8, color=color)
    text(slide, Inches(12.15), Inches(7.1), Inches(0.7), Inches(0.22), f"{n:02d}/{TOTAL:02d}", size=8, color=color, align=PP_ALIGN.RIGHT)


def _fit_size(iw: int, ih: int, max_w, max_h) -> tuple[int, int]:
    if iw <= 0 or ih <= 0:
        return int(max_w), int(max_h)
    scale = min(float(max_w) / iw, float(max_h) / ih)
    return int(iw * scale), int(ih * scale)


def add_media(
    slide,
    path: Path | None,
    l,
    t,
    max_w,
    max_h,
    *,
    caption: str | None = None,
    accent=C.mint,
    theme: str = "dark",
    frame: bool = True,
):
    fill = C.night_2 if theme == "dark" else C.white
    border = C.line if theme == "light" else BORDER_DARK

    if path is None or not path.exists():
        rect(slide, l, t, max_w, max_h, fill, border, radius=True)
        text(
            slide,
            l + Inches(0.2),
            t + max_h / 2 - Inches(0.2),
            max_w - Inches(0.4),
            Inches(0.4),
            "Скриншот: добавить позже",
            size=11,
            color=C.muted,
            align=PP_ALIGN.CENTER,
        )
        return None

    from PIL import Image

    with Image.open(path) as img:
        iw, ih = img.size

    final_w, final_h = _fit_size(iw, ih, max_w, max_h)
    final_l = l + int((max_w - final_w) / 2)
    final_t = t + int((max_h - final_h) / 2)
    if frame:
        rect(slide, l, t, max_w, max_h, fill, border, radius=True)
    pic = slide.shapes.add_picture(str(path), final_l, final_t, width=final_w, height=final_h)
    if caption:
        gen.add_caption(slide, l, t + max_h - Inches(0.18), max_w, caption, accent, theme=theme)
    return pic


def point_row(slide, y, num: str, head: str, accent, *, dark: bool = True, x=0.75, w=4.4):
    fg = C.white if dark else C.ink
    pill(slide, Inches(x), y, Inches(0.45), num, accent, C.ink if accent == C.lime else C.white)
    text(slide, Inches(x + 0.65), y + Inches(0.02), Inches(w), Inches(0.4), head, size=14, color=fg)
    return y + Inches(0.7)


# ─── slides ───────────────────────────────────────────────────────────────────


def slide_cover(prs, n):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, C.e8_black)
    rect(s, Inches(0.65), Inches(1.72), Inches(2.6), Pt(3), C.e8_coral)
    oval(s, Inches(11.95), Inches(0.5), Inches(0.75), Inches(0.75), C.e8_teal)
    rect(s, Inches(0), Inches(6.35), Inches(4.2), Inches(1.15), C.e8_teal)
    rect(s, Inches(4.2), Inches(6.35), W - Inches(4.2), Inches(1.15), C.e8_coral)

    add_logo(s, LOGO, Inches(0.65), Inches(0.5), Inches(2.4), Inches(1.1))
    text(s, Inches(0.65), Inches(2.05), Inches(4.8), Inches(1.4), "Datapipe", size=44, bold=True, color=C.white)
    text(
        s,
        Inches(0.65),
        Inches(3.55),
        Inches(4.6),
        Inches(1.0),
        "Инкрементальные ML-пайплайны\nна уровне каждой записи",
        size=15,
        color=C.e8_coral_muted,
    )
    text(s, Inches(0.65), Inches(5.0), Inches(4.4), Inches(0.4), "фреймворк · Ops UI · Data Collector", size=12, color=C.e8_teal_muted)

    add_media(s, PITCH / "graph-overview.png", Inches(5.4), Inches(1.2), Inches(7.3), Inches(4.7), caption="Datapipe Ops", accent=C.e8_teal, theme="dark", frame=False)
    text(s, Inches(0.75), Inches(6.55), Inches(12.0), Inches(0.45), "Вводная · учебный блок Epoch8", size=18, bold=True, color=C.e8_black, align=PP_ALIGN.CENTER)
    footer(s, n, dark=True)


def slide_problem(prs, n):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, C.night)
    title(s, "Проблема", "ML-пайплайны ломаются не на модели, а на данных", dark=True)

    items = [
        ("01", "Нет версионности", "Непонятно, на какой версии данных обучена модель и можно ли её воспроизвести.", C.coral),
        ("02", "Обвязка вместо инженерии", "Airflow, MLflow, Label Studio, S3 и скрипты живут порознь. Единой модели состояния нет.", C.mint),
        ("03", "Пересчёт всего заново", "Появились новые данные. Типичная реакция: прогнать пайплайн целиком «на всякий случай».", C.lime),
    ]
    x = Inches(0.65)
    w = Inches(3.95)
    for num, head, body, accent in items:
        rect(s, x, Inches(2.05), w, Inches(4.4), C.night_2, BORDER_DARK, radius=True)
        rect(s, x, Inches(2.05), w, Inches(0.1), accent)
        text(s, x + Inches(0.3), Inches(2.4), w - Inches(0.55), Inches(0.55), num, size=28, bold=True, color=accent)
        text(s, x + Inches(0.3), Inches(3.15), w - Inches(0.55), Inches(1.0), head, size=18, bold=True, color=C.white)
        text(s, x + Inches(0.3), Inches(4.35), w - Inches(0.55), Inches(1.7), body, size=13, color=MUTED_DARK)
        x += w + Inches(0.2)

    footer(s, n, dark=True)


def slide_what(prs, n):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, C.night)
    title(s, "Что такое Datapipe?", "Python-фреймворк для durable, incremental batch processing", dark=True)

    flow = ["images", "split", "train", "inference", "metrics", "fiftyone"]
    colors = [C.cobalt, C.mint, C.lime, C.coral, C.cobalt, C.mint]
    x = Inches(0.75)
    for label, color in zip(flow, colors):
        pill(s, x, Inches(2.05), Inches(1.8), label, color, C.ink if color == C.lime else C.white)
        x += Inches(2.05)

    points = [
        ("Record-level", "Изменилась строка. Пересчитываются только её downstream-шаги.", C.cobalt),
        ("Durable state", "Прерванный прогон продолжается с места остановки.", C.mint),
        ("Дешёвая итерация", "Не нужно перегонять всё. Можно падать и повторять безопасно.", C.lime),
        ("Curation UI", "Граф, runs, метрики, GT и prediction в одном интерфейсе.", C.coral),
    ]
    y = Inches(2.95)
    for head, body, accent in points:
        rect(s, Inches(0.75), y, Inches(11.85), Inches(0.85), C.night_2, BORDER_DARK, radius=True)
        rect(s, Inches(0.75), y, Inches(0.08), Inches(0.85), accent)
        text(s, Inches(1.1), y + Inches(0.12), Inches(3.2), Inches(0.55), head, size=14, bold=True, color=C.white, anchor=MSO_ANCHOR.MIDDLE)
        text(s, Inches(4.4), y + Inches(0.12), Inches(7.8), Inches(0.55), body, size=13, color=MUTED_DARK, anchor=MSO_ANCHOR.MIDDLE)
        y += Inches(0.98)

    footer(s, n, dark=True)


def slide_surfaces(prs, n):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, C.night)
    title(s, "Три поверхности, один пайплайн", "Код задаёт. Skills поднимают. UI наблюдает и запускает.", dark=True)

    surfaces = [
        ("01", "Python", "Catalog и Pipeline.\nПайплайн живёт в коде.", C.cobalt, TAGS / "python-catalog.png"),
        ("02", "Agentic Skills", "Env, сервисы, данные.\nАгент поднимает стенд.", C.mint, TAGS / "skills-setup.png"),
        ("03", "Datapipe Ops", "Граф, runs, metrics.\nЗапуск из живого UI.", C.lime, TAGS / "ui-retrain.png"),
    ]
    x = Inches(0.65)
    for num, head, body, accent, media in surfaces:
        rect(s, x, Inches(1.95), Inches(3.95), Inches(4.6), C.night_2, BORDER_DARK, radius=True)
        rect(s, x, Inches(1.95), Inches(0.08), Inches(4.6), accent)
        pill(s, x + Inches(0.25), Inches(2.15), Inches(0.55), num, accent, C.ink if accent == C.lime else C.white)
        text(s, x + Inches(0.95), Inches(2.2), Inches(2.7), Inches(0.35), head, size=16, bold=True, color=C.white)
        text(s, x + Inches(0.25), Inches(2.7), Inches(3.45), Inches(0.75), body, size=12, color=MUTED_DARK)
        add_media(s, media, x + Inches(0.2), Inches(3.6), Inches(3.55), Inches(2.7), theme="dark", frame=False, accent=accent)
        x += Inches(4.15)

    footer(s, n, dark=True)


def slide_python(prs, n):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    title(s, "Python: Catalog + Pipeline", "Декларативный граф: таблицы и трансформации")

    add_media(s, TAGS / "python-catalog.png", Inches(0.65), Inches(1.9), Inches(5.9), Inches(4.7), caption="Catalog", accent=C.cobalt, theme="light")
    add_media(s, TAGS / "python-pipeline.png", Inches(6.8), Inches(1.9), Inches(5.9), Inches(4.7), caption="Pipeline", accent=C.mint, theme="light")
    footer(s, n)


def slide_ops_graph(prs, n):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, C.night)
    title(s, "Datapipe Ops: граф", "Оранжевые узлы: данные. Синие: transforms и группы шагов.", dark=True)

    points = [
        ("1", "Структура пайплайна видна целиком", C.cobalt),
        ("2", "Статус этапов без чтения кода", C.mint),
        ("3", "Данные, разметка, обучение, метрики", C.lime),
    ]
    y = Inches(2.15)
    for num, head, accent in points:
        y = point_row(s, y, num, head, accent, dark=True, w=4.0)
        y += Inches(0.15)

    add_media(s, PITCH / "graph-overview.png", Inches(5.3), Inches(1.85), Inches(7.4), Inches(4.8), caption="Pipeline graph · train", accent=C.mint, theme="dark", frame=False)
    footer(s, n, dark=True)


def slide_runs(prs, n):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, C.night)
    title(s, "Каждый запуск с историей", "Статус, длительность, логи по шагам. Resume без ручной уборки.", dark=True)

    add_media(s, PITCH / "run-train-logs.png", Inches(0.55), Inches(1.85), Inches(6.15), Inches(4.8), caption="Train + логи эпох", accent=C.cobalt, theme="dark", frame=False)
    add_media(s, PITCH / "run-metrics-logs.png", Inches(6.9), Inches(1.85), Inches(5.9), Inches(4.8), caption="count-metrics", accent=C.mint, theme="dark", frame=False)
    footer(s, n, dark=True)


def slide_metrics(prs, n):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    title(s, "Метрики без чтения логов", "Precision, recall, F1 по модели, классу и тегу")

    points = [
        ("1", "Качество видно в UI", C.coral),
        ("2", "Разбивка по классам", C.cobalt),
        ("3", "Теги показывают слепые зоны", C.mint),
    ]
    y = Inches(2.2)
    for num, head, accent in points:
        y = point_row(s, y, num, head, accent, dark=False, w=4.0)
        y += Inches(0.25)

    add_media(s, PITCH / "metrics-table.png", Inches(5.2), Inches(1.9), Inches(7.5), Inches(4.7), caption="Metrics", accent=C.coral, theme="light")
    footer(s, n)


def slide_integrations(prs, n):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, C.night)
    title(s, "Интеграции", "Готовый набор для ML data work", dark=True)

    items = [
        ("Аннотации", "CVAT, Label Studio", C.coral),
        ("Визуализация", "FiftyOne", C.mint),
        ("Обучение", "YOLO, Ultralytics", C.lime),
        ("База данных", "PostgreSQL", C.cobalt),
        ("Хранилище", "MinIO, S3", C.coral),
        ("Оркестрация", "Airflow, K8s, CLI", C.mint),
    ]
    positions = [
        (0.75, 2.15),
        (4.75, 2.15),
        (8.75, 2.15),
        (0.75, 4.35),
        (4.75, 4.35),
        (8.75, 4.35),
    ]
    for (x, y), (head, body, accent) in zip(positions, items):
        rect(s, Inches(x), Inches(y), Inches(3.7), Inches(1.7), C.night_2, BORDER_DARK, radius=True)
        rect(s, Inches(x), Inches(y), Inches(0.08), Inches(1.7), accent)
        text(s, Inches(x + 0.3), Inches(y + 0.35), Inches(3.15), Inches(0.4), head, size=13, bold=True, color=MUTED_DARK)
        text(s, Inches(x + 0.3), Inches(y + 0.85), Inches(3.15), Inches(0.5), body, size=18, bold=True, color=C.white)

    footer(s, n, dark=True)


def slide_collector_bridge(prs, n):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, C.night)
    title(s, "Связь с Data Collector", "Как пакет из collector проходит через Datapipe", dark=True)

    steps = [
        ("1", "Commit пакета", "В collector session переходит в completed", C.cobalt),
        ("2", "Trigger", "Webhook или API вызывает пайплайн", C.mint),
        ("3", "Обработка", "Инференс, CVAT, запись результатов", C.lime),
        ("4", "Результат в админке", "Слои viz на пакете: inference, GT, CVAT", C.coral),
    ]
    y = Inches(2.1)
    for num, head, body, accent in steps:
        rect(s, Inches(0.75), y, Inches(11.85), Inches(1.05), C.night_2, BORDER_DARK, radius=True)
        rect(s, Inches(0.75), y, Inches(0.08), Inches(1.05), accent)
        pill(s, Inches(1.0), y + Inches(0.28), Inches(0.5), num, accent, C.ink if accent == C.lime else C.white)
        text(s, Inches(1.7), y + Inches(0.15), Inches(4.0), Inches(0.35), head, size=15, bold=True, color=C.white)
        text(s, Inches(5.9), y + Inches(0.28), Inches(6.3), Inches(0.5), body, size=13, color=MUTED_DARK, anchor=MSO_ANCHOR.MIDDLE)
        y += Inches(1.2)

    footer(s, n, dark=True)


def slide_video(prs, n):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, C.night)
    title(s, "Видео: Datapipe суть", "Что это, зачем, как устроен граф", dark=True)

    steps = [
        ("1", "Открыть Ops на учебном пайплайне", C.cobalt),
        ("2", "Показать граф: таблицы и transforms", C.mint),
        ("3", "Запустить шаг, посмотреть run и логи", C.lime),
        ("4", "Открыть metrics", C.coral),
        ("5", "Связать с collector: вход и результат", C.cobalt),
    ]
    y = Inches(2.0)
    for num, head, accent in steps:
        pill(s, Inches(0.8), y, Inches(0.5), num, accent, C.ink if accent == C.lime else C.white)
        text(s, Inches(1.5), y + Inches(0.02), Inches(4.5), Inches(0.4), head, size=13, color=C.white)
        rect(s, Inches(0.8), y + Inches(0.55), Inches(5.0), Pt(1), DIVIDER)
        y += Inches(0.75)

    text(s, Inches(0.8), Inches(6.2), Inches(5.0), Inches(0.4), "Файл: video/datapipe-overview.mp4", size=10, color=MUTED_DARK)

    rect(s, Inches(6.5), Inches(1.95), Inches(6.2), Inches(4.55), C.night_2, BORDER_DARK, radius=True)
    text(s, Inches(6.5), Inches(3.7), Inches(6.2), Inches(0.5), "▶  Вставить видео", size=20, bold=True, color=C.white, align=PP_ALIGN.CENTER)
    text(
        s,
        Inches(6.7),
        Inches(4.35),
        Inches(5.8),
        Inches(0.5),
        "PowerPoint: Вставка, Видео",
        size=12,
        color=MUTED_DARK,
        align=PP_ALIGN.CENTER,
    )
    footer(s, n, dark=True)


def main():
    prs = new_deck()
    n = 1
    slide_cover(prs, n); n += 1
    slide_problem(prs, n); n += 1
    slide_what(prs, n); n += 1
    slide_surfaces(prs, n); n += 1
    slide_python(prs, n); n += 1
    slide_ops_graph(prs, n); n += 1
    slide_runs(prs, n); n += 1
    slide_metrics(prs, n); n += 1
    slide_integrations(prs, n); n += 1
    slide_collector_bridge(prs, n); n += 1
    slide_video(prs, n)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"Wrote {OUT} ({TOTAL} slides)")


if __name__ == "__main__":
    main()
