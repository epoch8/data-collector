"""Диаграммы для презентации обучения."""

from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

BG = RGBColor(0x12, 0x12, 0x18)
PANEL = RGBColor(0x18, 0x20, 0x30)
CARD = RGBColor(0x1E, 0x2A, 0x3D)
TEAL = RGBColor(0x2D, 0xD4, 0xBF)
GOLD = RGBColor(0xD4, 0xA8, 0x4B)
GREEN = RGBColor(0x3D, 0xCC, 0x85)
PURPLE = RGBColor(0x9B, 0x7E, 0xD0)
BLUE = RGBColor(0x5B, 0x8D, 0xEF)
CORAL = RGBColor(0xE8, 0x7A, 0x7E)
TEXT = RGBColor(0xE8, 0xEA, 0xED)
MUTED = RGBColor(0x9C, 0xA3, 0xAF)
BORDER = RGBColor(0x3A, 0x45, 0x58)


def _canvas(slide, l, t, w, h):
    outer = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    outer.fill.solid()
    outer.fill.fore_color.rgb = BG
    outer.line.color.rgb = BORDER
    outer.line.width = Pt(1)
    outer.text_frame.clear()
    return outer


def _box(slide, l, t, w, h, fill, line, title, sub="", ts=10, ss=7):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.color.rgb = line
    s.line.width = Pt(1.5)
    tf = s.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Pt(6)
    tf.margin_right = Pt(6)
    p0 = tf.paragraphs[0]
    p0.text = title
    p0.font.size = Pt(ts)
    p0.font.bold = True
    p0.font.color.rgb = TEXT
    p0.alignment = PP_ALIGN.CENTER
    if sub:
        p1 = tf.add_paragraph()
        p1.text = sub
        p1.font.size = Pt(ss)
        p1.font.color.rgb = MUTED
        p1.alignment = PP_ALIGN.CENTER
        p1.space_before = Pt(2)
    return s


def _arrow_h(slide, x1, y, x2, color, dashed=False):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y, x2, y)
    c.line.color.rgb = color
    c.line.width = Pt(2)
    if dashed:
        c.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    return c


def _arrow_v(slide, x, y1, y2, color):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x, y1, x, y2)
    c.line.color.rgb = color
    c.line.width = Pt(2)
    return c


def _caption(slide, l, t, w, text, color=MUTED, size=8):
    box = slide.shapes.add_textbox(l, t, w, Inches(0.22))
    p = box.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.alignment = PP_ALIGN.CENTER


def draw_learning_pipeline(slide, l, t, w, h):
    """Путь обучения: от предметки до аттестации."""
    _canvas(slide, l, t, w, h)
    pad = Inches(0.12)
    ix, iy = l + pad, t + Inches(0.22)
    iw = w - pad * 2
    gap = Inches(0.08)
    n = 6
    bw = (iw - gap * (n - 1)) / n
    bh = Inches(0.95)
    y = iy + Inches(0.15)
    nodes = [
        ("Предметка\nбонитировка", GOLD),
        ("Съёмка\nв приложении", PURPLE),
        ("Upload\nи админка", BLUE),
        ("datapipe\nCVAT", TEAL),
        ("Эксплуатация\nинциденты", CORAL),
        ("Аттестация\nв проде", GREEN),
    ]
    for i, (label, col) in enumerate(nodes):
        x = ix + i * (bw + gap)
        _box(slide, x, y, bw, bh, CARD, col, label.split("\n")[0], label.split("\n")[1] if "\n" in label else "", 9, 6)
        if i < n - 1:
            _arrow_h(slide, x + bw, y + bh / 2, x + bw + gap, col)
    _caption(slide, ix, t + h - Inches(0.28), iw, "Участник проходит тот же путь, что и реальные данные Korovas", TEAL, 8)


def draw_three_phases(slide, l, t, w, h):
    """Три фазы программы вместо сухой таблицы блоков."""
    _canvas(slide, l, t, w, h)
    pad = Inches(0.14)
    ix, iy = l + pad, t + Inches(0.2)
    iw = w - pad * 2
    lane_h = Inches(1.35)
    gap = Inches(0.22)

    phases = [
        ("Фаза 1 · нед. 1–4", "Вход в продукт", "Предметка · инженерия · архитектура · первый E2E", GOLD),
        ("Фаза 2 · нед. 5–7", "Два трека → один результат", "Full-stack и ML параллельно на korovas-training", PURPLE),
        ("Фаза 3 · нед. 8–15", "Зрелость и передача", "Эксплуатация · прод · аттестация · бараны", GREEN),
    ]
    y = iy
    for title, head, body, col in phases:
        _box(slide, ix, y, iw, lane_h, CARD, col, title, head, 10, 8)
        box = slide.shapes.add_textbox(ix + Inches(0.15), y + Inches(0.62), iw - Inches(0.3), Inches(0.55))
        p = box.text_frame.paragraphs[0]
        p.text = body
        p.font.size = Pt(9)
        p.font.color.rgb = MUTED
        y += lane_h + gap


def draw_format_balance(slide, l, t, w, h):
    """Баланс форматов: видео мало, практика много."""
    _canvas(slide, l, t, w, h)
    pad = Inches(0.14)
    ix, iy = l + pad, t + Inches(0.25)
    iw = w - pad * 2
    _caption(slide, ix, iy - Inches(0.05), iw, "Соотношение форматов", TEXT, 9)

    items = [
        ("Микровидео", "35–45", 0.18, MUTED),
        ("Семинары", "20–24", 0.22, BLUE),
        ("Практикумы", "18–24", 0.55, TEAL),
        ("Инциденты", "8–12", 0.28, CORAL),
        ("Ролевые задачи", "8–10", 0.35, PURPLE),
        ("Ревью", "3–5", 0.15, GREEN),
    ]
    bar_h = Inches(0.38)
    gap = Inches(0.1)
    y = iy + Inches(0.2)
    max_w = iw * 0.72
    for name, count, frac, col in items:
        label_w = iw * 0.28
        slide.shapes.add_textbox(ix, y + Inches(0.06), label_w, bar_h).text_frame.paragraphs[0].text = name
        slide.shapes.add_textbox(ix, y + Inches(0.06), label_w, bar_h).text_frame.paragraphs[0].font.size = Pt(9)
        slide.shapes.add_textbox(ix, y + Inches(0.06), label_w, bar_h).text_frame.paragraphs[0].font.color.rgb = TEXT
        bx = ix + label_w
        bw = max_w * frac
        _box(slide, bx, y, bw, bar_h, col, col, count, "", 9, 7)
        y += bar_h + gap


def draw_stand_flow(slide, l, t, w, h):
    """Стенд: проекты и сценарии."""
    _canvas(slide, l, t, w, h)
    pad = Inches(0.14)
    ix, iy = l + pad, t + Inches(0.18)
    iw = w - pad * 2

    projects = [
        ("korovas-training", "основная практика", TEAL),
        ("korovas-broken", "инциденты", CORAL),
        ("demo-basic / demo-cv", "универсальные кейсы", BLUE),
    ]
    y = iy
    for name, role, col in projects:
        _box(slide, ix, y, iw, Inches(0.55), CARD, col, name, role, 10, 7)
        y += Inches(0.65)

    y += Inches(0.12)
    _caption(slide, ix, y, iw, "Сценарий успеха", TEAL, 8)
    y += Inches(0.22)
    steps = ["Съёмка", "Upload", "datapipe", "Viz", "Метрики"]
    sw = (iw - Inches(0.04 * 4)) / 5
    for i, step in enumerate(steps):
        x = ix + i * (sw + Inches(0.04))
        _box(slide, x, y, sw, Inches(0.48), PANEL, TEAL, step, "", 8, 7)
        if i < len(steps) - 1:
            _arrow_h(slide, x + sw, y + Inches(0.24), x + sw + Inches(0.04), TEAL)


def draw_parallel_weeks(slide, l, t, w, h):
    """Недели 5–7: FS и ML сходятся на E2E."""
    _canvas(slide, l, t, w, h)
    pad = Inches(0.14)
    ix, iy = l + pad, t + Inches(0.22)
    iw = w - pad * 2
    lh = Inches(0.9)
    gap = Inches(0.35)

    _box(slide, ix, iy, iw, lh, RGBColor(0x2A, 0x24, 0x38), PURPLE,
         "Full-stack", "config · upload · viz · админка", 10, 7)
    y2 = iy + lh + gap
    _box(slide, ix, y2, iw, lh, RGBColor(0x1A, 0x2E, 0x28), GREEN,
         "ML Engineer", "pipeline · CVAT · inference · метрики", 10, 7)
    merge_y = y2 + lh + Inches(0.28)
    _box(slide, ix + iw * 0.08, merge_y, iw * 0.84, Inches(0.62), CARD, TEAL,
         "Общий пакет korovas-training", "результат виден в админке", 10, 7)
    _arrow_v(slide, ix + iw / 2, iy + lh, iy + lh + gap / 2, PURPLE)
    _arrow_v(slide, ix + iw / 2, y2 + lh, merge_y, GREEN)
