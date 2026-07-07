"""
Generate customer-facing presentation for the Korovas training implementation.

Run:
    python docs/business/generate_training_presentation.py

Output:
    docs/business/Training-Program.pptx
"""

from __future__ import annotations

import math
from io import BytesIO
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

try:
    from PIL import Image, ImageDraw, ImageFilter
except ImportError:  # pragma: no cover - fallback for minimal envs
    Image = None
    ImageDraw = None
    ImageFilter = None

ROOT = Path(__file__).resolve().parent
REPO = ROOT.parent.parent
IMG = REPO / "specs" / "presentation" / "img"
LOGO = REPO / "e8-team-logo-1024.png"
OUT = ROOT / "Training-Program.pptx"

W = Inches(13.333)
H = Inches(7.5)
TOTAL = 7


class C:
    paper = RGBColor(0xFB, 0xFA, 0xF4)
    paper_2 = RGBColor(0xF1, 0xEE, 0xE6)
    ink = RGBColor(0x13, 0x16, 0x1D)
    muted = RGBColor(0x5F, 0x66, 0x73)
    faint = RGBColor(0xB6, 0xB8, 0xB2)
    line = RGBColor(0xDD, 0xDA, 0xCF)
    cobalt = RGBColor(0x22, 0x4E, 0xF2)
    lime = RGBColor(0xB7, 0xF2, 0x40)
    coral = RGBColor(0xFF, 0x6B, 0x57)
    mint = RGBColor(0x20, 0xC9, 0x9A)
    night = RGBColor(0x09, 0x0D, 0x16)
    night_2 = RGBColor(0x13, 0x1A, 0x2A)
    white = RGBColor(0xFF, 0xFF, 0xFF)
    e8_black = RGBColor(0x00, 0x00, 0x00)
    e8_coral = RGBColor(0xE8, 0x6B, 0x5E)
    e8_teal = RGBColor(0x3D, 0xB8, 0x9A)
    e8_coral_muted = RGBColor(0xC9, 0x8A, 0x82)
    e8_teal_muted = RGBColor(0x8A, 0xC9, 0xB5)


def new_deck() -> Presentation:
    prs = Presentation()
    prs.slide_width = W
    prs.slide_height = H
    return prs


def solid(shape, color: RGBColor) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = color


def no_line(shape) -> None:
    shape.line.fill.background()


def line(shape, color: RGBColor, width: float = 1.0) -> None:
    shape.line.color.rgb = color
    shape.line.width = Pt(width)


def rect(slide, l, t, w, h, fill: RGBColor, border: RGBColor | None = None, radius: bool = False):
    kind = MSO_SHAPE.ROUNDED_RECTANGLE if radius else MSO_SHAPE.RECTANGLE
    shape = slide.shapes.add_shape(kind, l, t, w, h)
    solid(shape, fill)
    if border:
        line(shape, border, 0.8)
    else:
        no_line(shape)
    return shape


def oval(slide, l, t, w, h, fill: RGBColor, border: RGBColor | None = None):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, l, t, w, h)
    solid(shape, fill)
    if border:
        line(shape, border, 0.8)
    else:
        no_line(shape)
    return shape


def text(
    slide,
    l,
    t,
    w,
    h,
    value: str,
    *,
    size=14,
    bold=False,
    color: RGBColor | None = None,
    align=PP_ALIGN.LEFT,
    anchor=MSO_ANCHOR.TOP,
):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.text = value
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color or C.ink
    p.alignment = align
    return box


def bg(slide, color: RGBColor = C.paper) -> None:
    solid(slide.background, color)


def footer(slide, n: int, dark: bool = False) -> None:
    color = RGBColor(0x7C, 0x83, 0x90) if dark else C.faint
    text(slide, Inches(0.65), Inches(7.1), Inches(4.5), Inches(0.22), "Epoch8 · обучение Korovas", size=8, color=color)
    text(
        slide,
        Inches(12.15),
        Inches(7.1),
        Inches(0.7),
        Inches(0.22),
        f"{n:02d}/{TOTAL:02d}",
        size=8,
        color=color,
        align=PP_ALIGN.RIGHT,
    )


def title(slide, title_text: str, subtitle: str | None = None, *, dark: bool = False):
    color = C.white if dark else C.ink
    muted = RGBColor(0xB0, 0xB7, 0xC5) if dark else C.muted
    text(slide, Inches(0.65), Inches(0.55), Inches(10.2), Inches(0.9), title_text, size=30, bold=True, color=color)
    rect(slide, Inches(0.65), Inches(1.48), Inches(1.8), Pt(5), C.lime if dark else C.cobalt)
    if subtitle:
        text(slide, Inches(0.65), Inches(1.65), Inches(10.2), Inches(0.45), subtitle, size=12, color=muted)


def pill(slide, l, t, w, label: str, fill: RGBColor, fg: RGBColor = C.ink):
    p = rect(slide, l, t, w, Inches(0.36), fill, radius=True)
    p.text_frame.clear()
    para = p.text_frame.paragraphs[0]
    para.text = label
    para.font.size = Pt(9)
    para.font.bold = True
    para.font.color.rgb = fg
    para.alignment = PP_ALIGN.CENTER
    p.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE


def card(slide, l, t, w, h, heading: str, body: str, accent: RGBColor = C.cobalt, dark: bool = False):
    fill = C.night_2 if dark else C.white
    border = RGBColor(0x2D, 0x36, 0x48) if dark else C.line
    muted = RGBColor(0xA7, 0xB0, 0xC0) if dark else C.muted
    rect(slide, l, t, w, h, fill, border, radius=True)
    rect(slide, l, t, Inches(0.08), h, accent)
    text(slide, l + Inches(0.22), t + Inches(0.18), w - Inches(0.35), Inches(0.35), heading, size=13, bold=True, color=accent)
    text(slide, l + Inches(0.22), t + Inches(0.58), w - Inches(0.35), h - Inches(0.7), body, size=10, color=muted)


def _image_path(path: Path, *, prefer_dark: bool = False) -> Path:
    if not prefer_dark:
        return path
    dark = path.parent / "dark" / path.name
    return dark if dark.exists() else path


def _cover_crop(img: Image.Image, target_w: int, target_h: int) -> Image.Image:
    img = img.convert("RGBA")
    src_w, src_h = img.size
    scale = max(target_w / src_w, target_h / src_h)
    resized = img.resize((math.ceil(src_w * scale), math.ceil(src_h * scale)), Image.Resampling.LANCZOS)
    left = (resized.width - target_w) // 2
    top = (resized.height - target_h) // 2
    return resized.crop((left, top, left + target_w, top + target_h))


def _rounded_mask(size: tuple[int, int], radius: int) -> Image.Image:
    mask = Image.new("L", size, 0)
    draw = ImageDraw.Draw(mask)
    draw.rounded_rectangle((0, 0, size[0] - 1, size[1] - 1), radius=radius, fill=255)
    return mask


def _apply_rounded(img: Image.Image, radius: int) -> Image.Image:
    mask = _rounded_mask(img.size, radius)
    out = Image.new("RGBA", img.size, (0, 0, 0, 0))
    out.paste(img, (0, 0), mask)
    return out


def _hex_rgb(color: RGBColor) -> tuple[int, int, int]:
    return color[0], color[1], color[2]


def _render_shadow(canvas: Image.Image, box: tuple[int, int, int, int], radius: int, blur: int = 10, alpha: int = 90) -> None:
    shadow = Image.new("RGBA", canvas.size, (0, 0, 0, 0))
    draw = ImageDraw.Draw(shadow)
    draw.rounded_rectangle(box, radius=radius, fill=(0, 0, 0, alpha))
    shadow = shadow.filter(ImageFilter.GaussianBlur(blur))
    canvas.alpha_composite(shadow)


def add_caption(
    slide,
    l,
    t,
    w,
    caption: str,
    accent: RGBColor,
    *,
    theme: str = "dark",
):
    """Floating caption badge overlapping the bottom edge of a screenshot."""
    cap_h = Inches(0.36)
    pad_x = Inches(0.24)
    est_w = Inches(0.095) * len(caption) + pad_x * 2
    cap_w = min(w - Inches(0.12), max(Inches(1.35), est_w))
    cap_l = l + int((w - cap_w) / 2)
    cap_t = t

    fill = RGBColor(0x10, 0x14, 0x1C) if theme == "dark" else C.white
    border = RGBColor(0x35, 0x3D, 0x4C) if theme == "dark" else C.line
    fg = C.white if theme == "dark" else C.ink

    rect(slide, cap_l, cap_t, cap_w, cap_h, fill, border, radius=True)
    rect(slide, cap_l, cap_t + Inches(0.08), Inches(0.07), cap_h - Inches(0.16), accent)
    text(
        slide,
        cap_l + Inches(0.18),
        cap_t + Inches(0.07),
        cap_w - Inches(0.24),
        cap_h - Inches(0.12),
        caption,
        size=9,
        bold=True,
        color=fg,
        align=PP_ALIGN.CENTER,
        anchor=MSO_ANCHOR.MIDDLE,
    )


def _render_panel(
    path: Path,
    width_px: int,
    height_px: int,
    *,
    style: str = "panel",
    accent: RGBColor = C.cobalt,
    theme: str = "dark",
) -> BytesIO | None:
    if Image is None or not path.exists():
        return None

    pad = 16
    outer_w = width_px
    outer_h = height_px
    inner_w = outer_w - pad * 2
    inner_h = outer_h - pad * 2
    if inner_w <= 0 or inner_h <= 0:
        return None

    accent_rgb = _hex_rgb(accent)
    canvas = Image.new("RGBA", (outer_w, outer_h), (0, 0, 0, 0))
    shell = (pad - 4, pad - 4, outer_w - pad + 4, outer_h - pad + 4)
    _render_shadow(canvas, shell, radius=22, alpha=55 if theme == "light" else 90)

    shell_img = Image.new("RGBA", (outer_w, outer_h), (0, 0, 0, 0))
    shell_draw = ImageDraw.Draw(shell_img)
    if theme == "light":
        shell_draw.rounded_rectangle(shell, radius=20, fill=(255, 255, 255, 255), outline=(*accent_rgb, 200), width=2)
    else:
        shell_draw.rounded_rectangle(shell, radius=20, fill=(19, 26, 42, 255), outline=(*accent_rgb, 180), width=2)
    canvas.alpha_composite(shell_img)

    with Image.open(path) as src:
        if style == "phone":
            bezel = 10
            phone_w = inner_w - bezel * 2
            phone_h = inner_h - bezel * 2
            shot = _cover_crop(src, phone_w, phone_h)
            shot = _apply_rounded(shot, 18)
            phone = Image.new("RGBA", (inner_w, inner_h), (0, 0, 0, 0))
            frame_draw = ImageDraw.Draw(phone)
            frame_draw.rounded_rectangle(
                (0, 0, inner_w - 1, inner_h - 1),
                radius=24,
                fill=(10, 12, 18, 255),
                outline=(70, 82, 104, 255),
                width=2,
            )
            notch_w = max(42, inner_w // 4)
            frame_draw.rounded_rectangle(
                ((inner_w - notch_w) // 2, 8, (inner_w + notch_w) // 2, 18),
                radius=4,
                fill=(6, 8, 12, 255),
            )
            phone.paste(shot, (bezel, bezel), shot)
            shot = phone
        else:
            title_h = 30
            inset = 8
            content_w = inner_w - inset * 2
            content_h = inner_h - title_h - inset
            shot = _cover_crop(src, content_w, content_h)
            shot = _apply_rounded(shot, 10)
            browser = Image.new("RGBA", (inner_w, inner_h), (0, 0, 0, 0))
            frame_draw = ImageDraw.Draw(browser)
            if theme == "light":
                frame_fill = (248, 247, 244, 255)
                chrome_fill = (236, 234, 228, 255)
                outline = (221, 218, 207, 255)
            else:
                frame_fill = (12, 16, 24, 255)
                chrome_fill = (24, 30, 42, 255)
                outline = (58, 68, 86, 255)
            frame_draw.rounded_rectangle(
                (0, 0, inner_w - 1, inner_h - 1),
                radius=16,
                fill=frame_fill,
                outline=outline,
                width=2,
            )
            frame_draw.rectangle((0, 0, inner_w - 1, title_h), fill=chrome_fill)
            frame_draw.ellipse((12, 10, 22, 20), fill=(232, 107, 94, 255))
            frame_draw.ellipse((28, 10, 38, 20), fill=(247, 198, 94, 255))
            frame_draw.ellipse((44, 10, 54, 20), fill=(61, 184, 154, 255))
            browser.paste(shot, (inset, title_h), shot)
            shot = browser

    canvas.paste(shot, (pad, pad), shot)

    buf = BytesIO()
    canvas.save(buf, format="PNG")
    buf.seek(0)
    return buf


def add_image_panel(
    slide,
    path: Path,
    l,
    t,
    w,
    h,
    *,
    style: str = "panel",
    accent: RGBColor = C.cobalt,
    caption: str | None = None,
    prefer_dark: bool = False,
    rotation: float = 0,
    theme: str = "dark",
):
    """Insert a polished screenshot card with crop-to-fill and device chrome."""
    resolved = _image_path(path, prefer_dark=prefer_dark)
    width_px = max(320, int(w / 914400 * 160))
    height_px = max(220, int(h / 914400 * 160))
    rendered = _render_panel(resolved, width_px, height_px, style=style, accent=accent, theme=theme)
    if rendered is None:
        frame = rect(slide, l, t, w, h, C.night_2, C.line, radius=True)
        frame.text_frame.clear()
        text(slide, l, t + h / 2 - Inches(0.15), w, Inches(0.3), resolved.name, size=9, color=C.faint, align=PP_ALIGN.CENTER)
        return frame

    pic = slide.shapes.add_picture(rendered, l, t, width=w, height=h)
    if rotation:
        pic.rotation = rotation
    if caption:
        add_caption(slide, l, t + h - Inches(0.2), w, caption, accent, theme=theme)
    return pic


def add_image_fit(slide, path: Path, l, t, w, h, *, bg_fill: RGBColor = C.white, caption: str | None = None, style: str = "browser"):
    """Backward-compatible wrapper around the polished screenshot panel."""
    accent = C.cobalt if bg_fill == C.white else C.e8_teal
    dark = bg_fill != C.white
    return add_image_panel(
        slide,
        path,
        l,
        t,
        w,
        h,
        style=style,
        accent=accent,
        caption=caption,
        prefer_dark=dark,
    )


def arrow(slide, x1, y1, x2, y2, color=C.cobalt):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    line(c, color, 2)
    return c


def add_logo(slide, path: Path, l, t, max_w, max_h, *, align_left: bool = True):
    """Place logo preserving aspect ratio, no frame."""
    if not path.exists():
        text(slide, l, t, max_w, max_h, "E8", size=48, bold=True, color=C.e8_coral)
        return None
    if Image is None:
        return slide.shapes.add_picture(str(path), l, t, width=max_w)
    with Image.open(path) as img:
        img_w, img_h = img.size
    box_ratio = float(max_w) / float(max_h)
    img_ratio = img_w / img_h
    if img_ratio > box_ratio:
        final_w = max_w
        final_h = int(max_w / img_ratio)
    else:
        final_h = max_h
        final_w = int(max_h * img_ratio)
    final_l = l if align_left else l + int((max_w - final_w) / 2)
    final_t = t + int((max_h - final_h) / 2)
    return slide.shapes.add_picture(str(path), final_l, final_t, width=final_w, height=final_h)


def slide_cover(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, C.e8_black)

    # Brand accents — coral outline motif (E) + teal fill motif (8)
    rect(s, Inches(0.65), Inches(1.72), Inches(2.6), Pt(3), C.e8_coral)
    oval(s, Inches(11.95), Inches(0.5), Inches(0.75), Inches(0.75), C.e8_teal)
    rect(s, Inches(0), Inches(6.35), Inches(4.2), Inches(1.15), C.e8_teal)
    rect(s, Inches(4.2), Inches(6.35), W - Inches(4.2), Inches(1.15), C.e8_coral)

    add_logo(s, LOGO, Inches(0.65), Inches(0.5), Inches(2.4), Inches(1.1))

    text(
        s,
        Inches(0.65),
        Inches(2.05),
        Inches(4.8),
        Inches(2.5),
        "Передача\nсопровождения\nсистемы",
        size=38,
        bold=True,
        color=C.white,
    )
    text(
        s,
        Inches(0.65),
        Inches(4.75),
        Inches(4.5),
        Inches(0.7),
        "Программа обучения для full-stack и ML-инженера",
        size=13,
        color=C.e8_coral_muted,
    )
    text(
        s,
        Inches(0.65),
        Inches(5.3),
        Inches(4.5),
        Inches(0.45),
        "Korovas · data-collector · datapipe",
        size=11,
        color=C.e8_teal_muted,
    )

    text(
        s,
        Inches(5.55),
        Inches(0.75),
        Inches(7.0),
        Inches(0.4),
        "data-collector · datapipe · CVAT · инференс · эксплуатация",
        size=12,
        color=RGBColor(0x7A, 0x82, 0x8C),
        align=PP_ALIGN.RIGHT,
    )
    add_image_panel(
        s,
        IMG / "Flutter_en" / "1.jpg",
        Inches(5.35),
        Inches(1.55),
        Inches(2.05),
        Inches(3.75),
        style="phone",
        accent=C.e8_coral,
        caption="Съёмка в поле",
        prefer_dark=True,
        rotation=-2,
    )
    add_image_panel(
        s,
        IMG / "Config_en" / "1.png",
        Inches(10.35),
        Inches(1.65),
        Inches(2.35),
        Inches(3.55),
        style="browser",
        accent=C.e8_coral,
        caption="Конфигурация проекта",
        prefer_dark=True,
        rotation=2,
    )
    add_image_panel(
        s,
        IMG / "UI_en" / "5.png",
        Inches(7.15),
        Inches(1.25),
        Inches(3.35),
        Inches(4.05),
        style="browser",
        accent=C.e8_teal,
        caption="Визуализация инференса",
        prefer_dark=True,
    )

    text(
        s,
        Inches(0.75),
        Inches(6.55),
        Inches(12.0),
        Inches(0.45),
        "300 ч · 15 недель · 2 роли · учебный стенд",
        size=18,
        bold=True,
        color=C.e8_black,
        align=PP_ALIGN.CENTER,
    )
    text(s, Inches(0.65), Inches(7.08), Inches(4.5), Inches(0.22), "Epoch8 · обучение Korovas", size=8, color=RGBColor(0x55, 0x55, 0x55))
    text(s, Inches(12.15), Inches(7.08), Inches(0.7), Inches(0.22), "01/07", size=8, color=RGBColor(0x55, 0x55, 0x55), align=PP_ALIGN.RIGHT)


def slide_system(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    title(s, "Из чего состоит система", "И почему обучение строится вокруг пути одного пакета")
    cards = [
        ("data-collector", "мобильная съёмка, загрузка пакетов, админка, визуализация", C.cobalt),
        ("datapipe", "пайплайны обработки, CVAT, инференс, метрики", C.mint),
        ("слой Korovas", "бонитировка КРС, правила съёмки, экспертная оценка", C.coral),
    ]
    x = Inches(0.75)
    for head, body, color in cards:
        card(s, x, Inches(2.05), Inches(3.7), Inches(1.35), head, body, color)
        x += Inches(4.05)

    stages = [
        ("01", "Съёмка", "приложение", C.cobalt),
        ("02", "Загрузка", "пакет", C.coral),
        ("03", "Пайплайн", "datapipe", C.mint),
        ("04", "Просмотр", "визуализация и измерения", C.lime),
        ("05", "Протокол", "Монолит и внешние сервисы", C.coral),
    ]
    x = Inches(0.75)
    y = Inches(4.55)
    w = Inches(2.05)
    for i, (num, head, body, color) in enumerate(stages):
        rect(s, x, y, w, Inches(1.25), C.white, C.line, radius=True)
        pill(s, x + Inches(0.18), y + Inches(0.18), Inches(0.62), num, color, C.ink if color == C.lime else C.white)
        text(s, x + Inches(0.18), y + Inches(0.62), w - Inches(0.36), Inches(0.28), head, size=14, bold=True)
        text(s, x + Inches(0.18), y + Inches(0.9), w - Inches(0.36), Inches(0.3), body, size=7, color=C.muted)
        if i < len(stages) - 1:
            arrow(s, x + w, y + Inches(0.62), x + w + Inches(0.2), y + Inches(0.62), C.ink)
        x += w + Inches(0.24)
    footer(s, 2)


def slide_outcome(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s, C.night)
    title(s, "Финальный результат", "Что должна уметь команда после 15 недель", dark=True)
    outcomes = [
        ("01", "Провести пакет end-to-end", "От съёмки коровы до визуализации инференса в админке."),
        ("02", "Диагностировать сбой", "Загрузка, хранилище, datapipe, CVAT, пустая визуализация."),
        ("03", "Оценить качество ML", "Сравнить модель, GT и экспертную оценку по метрикам."),
        ("04", "Перенести подход", "Повторить сценарий на смежной предметной области, например баранах."),
    ]
    y = Inches(1.95)
    for num, head, body in outcomes:
        text(s, Inches(0.8), y, Inches(0.7), Inches(0.42), num, size=18, bold=True, color=C.lime)
        text(s, Inches(1.65), y, Inches(4.3), Inches(0.35), head, size=15, bold=True, color=C.white)
        text(s, Inches(1.65), y + Inches(0.38), Inches(5.5), Inches(0.35), body, size=10, color=RGBColor(0xAE, 0xB7, 0xC7))
        rect(s, Inches(0.8), y + Inches(0.88), Inches(5.5), Pt(1), RGBColor(0x2B, 0x34, 0x45))
        y += Inches(1.15)
    add_image_panel(
        s,
        IMG / "UI_en" / "5.png",
        Inches(7.15),
        Inches(1.55),
        Inches(5.35),
        Inches(4.65),
        style="browser",
        accent=C.lime,
        caption="Визуализация инференса",
        prefer_dark=True,
    )
    footer(s, 3, dark=True)


def slide_program(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    title(s, "15 недель как три этапа", "Внутри остаются 8 блоков roadmap, но маршрут для запуска проще")
    acts = [
        ("I", "Войти в систему", "нед. 1–4", "Предметная область, инженерная база, архитектура, первый E2E на стенде.", C.cobalt),
        ("II", "Разойтись по ролям", "нед. 5–7", "Full-stack работает с загрузкой и визуализацией, ML — с пайплайном, CVAT и метриками.", C.mint),
        ("III", "Принять сопровождение", "нед. 8–15", "Инциденты, production-практика, аттестация, перенос на новую предметку.", C.coral),
    ]
    x = Inches(0.75)
    for num, head, weeks, body, color in acts:
        rect(s, x, Inches(2.1), Inches(3.75), Inches(3.8), C.white, C.line, radius=True)
        oval(s, x + Inches(0.25), Inches(2.4), Inches(0.58), Inches(0.58), color)
        text(s, x + Inches(0.25), Inches(2.52), Inches(0.58), Inches(0.25), num, size=14, bold=True, color=C.white, align=PP_ALIGN.CENTER)
        text(s, x + Inches(0.25), Inches(3.15), Inches(3.1), Inches(0.5), head, size=22, bold=True)
        text(s, x + Inches(0.25), Inches(3.85), Inches(3.1), Inches(0.3), weeks, size=12, bold=True, color=C.muted)
        text(s, x + Inches(0.25), Inches(4.45), Inches(3.15), Inches(1.0), body, size=11, color=C.muted)
        x += Inches(4.05)
    footer(s, 4)


def slide_learning_process(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    title(s, "Способы передачи знаний", "За счёт чего команда погружается в проект и учится сопровождать систему")
    steps = [
        ("01", "Короткие видео", "дают контекст и вводный разбор"),
        ("02", "Семинары", "показываем рабочий процесс на стенде"),
        ("03", "Практические задания", "участник повторяет сценарий на своих руках"),
        ("04", "Разборы инцидентов", "учимся находить причину сбоя по логам и данным"),
        ("05", "Шаблоны и инструкции", "чек-листы, runbook, структура техотчёта"),
        ("06", "Ревью результата", "обратная связь и приёмка выполненного задания"),
    ]
    x = Inches(0.75)
    y = Inches(2.1)
    w = Inches(3.7)
    for i, (num, head, body) in enumerate(steps):
        col = i % 3
        row = i // 3
        sx = x + col * Inches(4.05)
        sy = y + row * Inches(1.65)
        color = [C.cobalt, C.mint, C.coral][col]
        rect(s, sx, sy, w, Inches(1.25), C.white, C.line, radius=True)
        pill(s, sx + Inches(0.22), sy + Inches(0.2), Inches(0.62), num, color, C.white)
        text(s, sx + Inches(1.0), sy + Inches(0.18), Inches(2.35), Inches(0.28), head, size=14, bold=True)
        text(s, sx + Inches(1.0), sy + Inches(0.55), Inches(2.35), Inches(0.5), body, size=9, color=C.muted)
    text(s, Inches(0.9), Inches(5.8), Inches(11.3), Inches(0.5), "Форматы чередуются: сначала короткое объяснение, затем демонстрация, практика и проверка результата.", size=15, bold=True, color=C.cobalt, align=PP_ALIGN.CENTER)
    footer(s, 5)


def slide_practice_base(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    title(s, "Практическая база и специфика", "Готовим демо-примеры, локальные стенды и поломанные кейсы")
    left_items = [
        ("korovas-datapipe", "формирование выборок, оценка метрик, анализ ошибок и работа с пайплайнами", C.mint),
        ("korovas-datacollector", "мобильное приложение, админка, загрузка пакетов и связь с datapipe", C.cobalt),
        ("korovas-broken", "набор намеренно сломанных кейсов: загрузка, datapipe, визуализация, хранилище", C.coral),
        ("локальные стенды", "шаблоны конфигураций для локального развёртывания БД, S3 и других сервисов", C.lime),
    ]
    y = Inches(2.0)
    for head, body, color in left_items:
        card(s, Inches(0.7), y, Inches(5.75), Inches(0.86), head, body, color)
        y += Inches(0.98)
    add_image_panel(
        s,
        IMG / "packege_list_slide_2.png",
        Inches(6.75),
        Inches(1.85),
        Inches(5.85),
        Inches(4.55),
        style="browser",
        accent=C.cobalt,
        caption="Список пакетов",
        prefer_dark=False,
        theme="light",
    )
    footer(s, 6)


def slide_final_roadmap(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    bg(s)
    rect(s, Inches(0), Inches(0), W, Inches(1.25), C.night)
    text(s, Inches(0.75), Inches(0.38), Inches(9.0), Inches(0.5), "Финальный roadmap подготовки", size=28, bold=True, color=C.white)
    items = [
        ("01", "Собрать учебные данные", "korovas-training, korovas-broken, примеры для CVAT и инференса"),
        ("02", "Подготовить конфигурации", "шаблоны для локального развёртывания БД, S3, data-collector и datapipe"),
        ("03", "Подготовить стартовые материалы", "микровидео, инструкции, чек-листы, шаблоны техотчётов"),
        ("04", "Описать практикумы", "E2E, загрузка, визуализация, пайплайн, CVAT, метрики, инциденты"),
        ("05", "Подготовить задания по ролям", "full-stack и ML-инженер, критерии ревью и приёмки"),
        ("06", "Запустить неделю 1", "демо полного сценария на коровах и первый разбор пути пакета"),
    ]
    y = Inches(1.75)
    for num, head, body in items:
        text(s, Inches(0.9), y, Inches(0.65), Inches(0.35), num, size=13, bold=True, color=C.cobalt)
        text(s, Inches(1.65), y, Inches(3.8), Inches(0.35), head, size=16, bold=True, color=C.ink)
        text(s, Inches(5.65), y + Inches(0.03), Inches(6.7), Inches(0.32), body, size=10, color=C.muted)
        rect(s, Inches(1.65), y + Inches(0.5), Inches(10.7), Pt(1), C.line)
        y += Inches(0.72)
    footer(s, 7)


def build() -> Path:
    prs = new_deck()
    slide_cover(prs)
    slide_system(prs)
    slide_outcome(prs)
    slide_program(prs)
    slide_learning_process(prs)
    slide_practice_base(prs)
    slide_final_roadmap(prs)
    prs.save(OUT)
    return OUT


if __name__ == "__main__":
    print(f"Saved: {build()}")
