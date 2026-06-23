"""Minimal vector diagrams for Data Collector presentation."""

from pptx.dml.color import RGBColor
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

BG_DEEP = RGBColor(0x0A, 0x0E, 0x18)
BG_CARD = RGBColor(0x1E, 0x2A, 0x3D)
BG_PANEL = RGBColor(0x14, 0x1A, 0x26)
ACCENT = RGBColor(0x5B, 0x8D, 0xEF)
ACCENT_GREEN = RGBColor(0x3D, 0xCC, 0x85)
ACCENT_GOLD = RGBColor(0xD4, 0xA8, 0x4B)
ACCENT_PURPLE = RGBColor(0x9B, 0x7E, 0xD0)
ACCENT_CYAN = RGBColor(0x8B, 0xC9, 0xF7)
TEXT = RGBColor(0xE8, 0xEA, 0xED)
TEXT_MUTED = RGBColor(0x9C, 0xA3, 0xAF)
BORDER = RGBColor(0x3A, 0x45, 0x58)


def _box(slide, l, t, w, h, fill, line, title, sub="", ts=11, ss=8):
    s = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    s.line.color.rgb = line
    s.line.width = Pt(1.5)
    tf = s.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Pt(8)
    tf.margin_right = Pt(8)
    p0 = tf.paragraphs[0]
    p0.text = title
    p0.font.size = Pt(ts)
    p0.font.bold = True
    p0.font.color.rgb = TEXT
    p0.alignment = PP_ALIGN.CENTER
    if sub:
        p1 = tf.add_paragraph()
        p1.text = sub
        p1.font.size = Pt(ss)
        p1.font.color.rgb = TEXT_MUTED
        p1.alignment = PP_ALIGN.CENTER
        p1.space_before = Pt(3)
    return s


def _arrow_h(slide, x1, y, x2, color=ACCENT, dashed=False):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y, x2, y)
    c.line.color.rgb = color
    c.line.width = Pt(2)
    if dashed:
        c.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    return c


def _arrow_v(slide, x, y1, y2, color=ACCENT, dashed=False):
    c = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x, y1, x, y2)
    c.line.color.rgb = color
    c.line.width = Pt(2)
    if dashed:
        c.line.dash_style = MSO_LINE_DASH_STYLE.DASH
    return c


def _caption(slide, l, t, w, text, color=TEXT_MUTED, size=8):
    box = slide.shapes.add_textbox(l, t, w, Inches(0.22))
    p = box.text_frame.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.alignment = PP_ALIGN.CENTER


def _canvas(slide, l, t, w, h):
    outer = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h)
    outer.fill.solid()
    outer.fill.fore_color.rgb = BG_DEEP
    outer.line.color.rgb = BORDER
    outer.line.width = Pt(1)
    outer.text_frame.clear()
    return outer


# ── Cover: Admin → Server → App + Admin ───────────────────────────────────────

def draw_cover_architecture(slide, l, t, w, h):
    _canvas(slide, l, t, w, h)
    pad = Inches(0.2)
    bw = w - pad * 2
    bh = Inches(0.52)
    gap = Inches(0.38)
    y = t + Inches(0.25)
    cx = l + w / 2

    nodes = [
        ("Staff Admin", "конфиг", RGBColor(0x3D, 0x34, 0x20), ACCENT_GOLD),
        ("Django Server", "API + хранилище", BG_CARD, ACCENT),
        ("Flutter App", "сбор", RGBColor(0x2A, 0x24, 0x38), ACCENT_PURPLE),
        ("Admin UI", "пакеты", RGBColor(0x1A, 0x2E, 0x28), ACCENT_GREEN),
    ]
    for title, sub, fill, line in nodes:
        _box(slide, l + pad, y, bw, bh, fill, line, title, sub, 10, 7)
        y += bh + gap
        if title != nodes[-1][0]:
            _arrow_v(slide, cx, y - gap + Inches(0.02), y - Inches(0.02), line)

    _arrow_h(slide, l + pad + bw * 0.55, t + Inches(0.25) + bh * 2 + gap * 2 + bh * 0.5,
              l + pad + bw * 0.85, ACCENT_PURPLE, dashed=True)


# ── Stack (legacy) ───────────────────────────────────────────────────────────

def draw_stack_diagram(slide, l, t, w, h):
    draw_product_architecture(slide, l, t, w, h)


# ── Product: who talks to whom (not a tech list) ───────────────────────────

def draw_product_architecture(slide, l, t, w, h):
    """Vertical interaction map: config down to server, data up from app to admins."""
    _canvas(slide, l, t, w, h)
    pad = Inches(0.14)
    ix, iy = l + pad, t + Inches(0.2)
    iw = w - pad * 2
    cx = ix + iw / 2
    bh = Inches(0.56)
    gap = Inches(0.22)

    _caption(slide, ix, iy - Inches(0.02), iw,
             "Кто с кем взаимодействует — не список технологий", ACCENT_CYAN, 9)

    # Staff → Django
    y = iy + Inches(0.18)
    staff_w = iw * 0.62
    _box(slide, cx - staff_w / 2, y, staff_w, bh,
         RGBColor(0x3D, 0x34, 0x20), ACCENT_GOLD, "Staff Admin", "настройка проекта и конфига", 10, 7)

    y2 = y + bh + gap
    _caption(slide, cx - Inches(0.55), y + bh, Inches(1.1), "конфиг\n+ права", ACCENT_GOLD, 7)
    _arrow_v(slide, cx, y + bh, y2, ACCENT_GOLD)

    srv_w = iw * 0.78
    _box(slide, cx - srv_w / 2, y2, srv_w, bh + Inches(0.08), BG_CARD, ACCENT,
         "Django + хранилище", "API · БД · медиа · синк из Git", 10, 7)

    # Git sidecar
    git_w = iw * 0.28
    git_x = ix + iw - git_w
    git_y = y2 + Inches(0.06)
    _box(slide, git_x, git_y, git_w, bh, RGBColor(0x1E, 0x24, 0x33), ACCENT_CYAN,
         "Git", "config · viz", 9, 7)
    _arrow_h(slide, git_x, git_y + bh / 2, cx + srv_w / 2, ACCENT_CYAN, dashed=True)
    _caption(slide, git_x - Inches(0.05), git_y + bh + Inches(0.02), git_w, "версии", ACCENT_CYAN, 7)

    y3 = y2 + bh + Inches(0.08) + gap
    half_w = (iw - Inches(0.12)) / 2

    _box(slide, ix, y3, half_w, bh, RGBColor(0x2A, 0x24, 0x38), ACCENT_PURPLE,
         "Flutter App", "сборщик в поле", 10, 7)
    _box(slide, ix + half_w + Inches(0.12), y3, half_w, bh, RGBColor(0x1A, 0x2E, 0x28), ACCENT_GREEN,
         "Client Admin", "просмотр пакетов", 10, 7)

    mid_y = y2 + bh + Inches(0.08) + gap / 2
    _arrow_v(slide, cx - half_w / 2 - Inches(0.06), y2 + bh + Inches(0.08), y3,
             ACCENT_PURPLE, dashed=False)
    _caption(slide, ix, mid_y, half_w, "синк конфига", ACCENT_PURPLE, 7)

    _arrow_v(slide, cx + half_w / 2 + Inches(0.06), y2 + bh + Inches(0.08), y3,
             ACCENT_GREEN, dashed=False)
    _caption(slide, ix + half_w + Inches(0.12), mid_y, half_w, "список пакетов", ACCENT_GREEN, 7)

    y4 = y3 + bh + Inches(0.14)
    _caption(slide, ix, y3 + bh, half_w, "upload пакета", ACCENT_PURPLE, 7)
    _arrow_v(slide, cx - half_w / 2 - Inches(0.06), y3 + bh, y4 - Inches(0.02), ACCENT_PURPLE)

    foot_h = min(Inches(0.52), t + h - y4 - Inches(0.08))
    if foot_h > Inches(0.3):
        _box(slide, ix, y4, iw, foot_h, RGBColor(0x14, 0x20, 0x18), ACCENT_GREEN,
             "Проект × N", "1 проект = 1 конфиг в Git · свои пакеты", 10, 7)


# ── Dual path: config lifecycle vs package lifecycle ─────────────────────────

def draw_dual_path_flow(slide, l, t, w, h):
    """Two lanes: config (setup → sync → viz) and package (collect → upload → view)."""
    _canvas(slide, l, t, w, h)
    pad = Inches(0.14)
    ix, iy = l + pad, t + Inches(0.18)
    iw = w - pad * 2
    lane_h = Inches(1.02)
    lane_gap = Inches(0.38)
    n = 3
    gap = Inches(0.1)
    bw = (iw - gap * (n - 1)) / n

    _caption(slide, ix, iy - Inches(0.04), iw,
             "Два пути: жизненный цикл конфига и жизненный цикл пакета данных", TEXT, 9)

    lanes = [
        ("Конфиг (config.json + viz.json в Git)", ACCENT_GOLD, [
            ("①", "Настройка\nв Staff Admin", "fields · flow · Git"),
            ("②", "Синк\nв приложение", "тот же сценарий в МП"),
            ("③", "Визуализации\nв админке", "viz.json → UI просмотра"),
        ]),
        ("Пакет данных (результат сбора)", ACCENT_PURPLE, [
            ("①", "Сбор\nв приложении", "сборщик по flow"),
            ("②", "Upload\nна сервер", "поля + фото + манифест"),
            ("③", "Просмотр\nв админке", "фильтры · workspace"),
        ]),
    ]

    y = iy + Inches(0.2)
    for lane_title, lane_color, steps in lanes:
        _caption(slide, ix, y, iw, lane_title, lane_color, 8)
        y += Inches(0.22)
        cy = y + lane_h / 2
        for i, (num, label, sub) in enumerate(steps):
            x = ix + i * (bw + gap)
            _box(slide, x, y, bw, lane_h, BG_CARD, lane_color, f"{num}  {label}", sub, 10, 7)
            if i < n - 1:
                _arrow_h(slide, x + bw, cy, x + bw + gap, lane_color)
        y += lane_h + lane_gap


# ── 5-step product flow ───────────────────────────────────────────────────────

def draw_five_step_flow(slide, l, t, w, h):
    """Horizontal pipeline: config → sync → collect → upload → view."""
    _canvas(slide, l, t, w, h)
    pad = Inches(0.14)
    ix = l + pad
    iw = w - pad * 2
    n = 5
    gap = Inches(0.12)
    bw = (iw - gap * (n - 1)) / n
    bh = min(Inches(1.05), h - Inches(0.28))
    iy = t + (h - bh) / 2

    steps = [
        ("①", "Конфиг", ACCENT_GOLD),
        ("②", "Синк в МП", ACCENT),
        ("③", "Сбор", ACCENT_PURPLE),
        ("④", "Upload", ACCENT_PURPLE),
        ("⑤", "Админка", ACCENT_GREEN),
    ]
    for i, (num, label, c) in enumerate(steps):
        x = ix + i * (bw + gap)
        _box(slide, x, iy, bw, bh, BG_CARD, c, num, label, 16, 10)
        if i < n - 1:
            _arrow_h(slide, x + bw, iy + bh / 2, x + bw + gap, c)


# ── Config: fields · flow · viz (with screen semantics) ─────────────────────

def draw_config_entities(slide, l, t, w, h):
    _canvas(slide, l, t, w, h)
    pad = Inches(0.16)
    ix, iy = l + pad, t + Inches(0.16)
    iw = w - pad * 2

    banner_h = Inches(0.42)
    banner = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, ix, iy, iw, banner_h)
    banner.fill.solid()
    banner.fill.fore_color.rgb = RGBColor(0x28, 0x32, 0x48)
    banner.line.color.rgb = ACCENT_GOLD
    banner.line.width = Pt(1)
    tf = banner.text_frame
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = tf.paragraphs[0]
    p.text = "Конфиг в Git состоит из этих частей:"
    p.font.size = Pt(10)
    p.font.bold = True
    p.font.color.rgb = ACCENT_GOLD
    p.alignment = PP_ALIGN.CENTER

    gap = Inches(0.1)
    y = iy + banner_h + Inches(0.14)
    bw = iw

    bh = Inches(0.68)
    hint_h = Inches(0.16)
    parts = [
        ("fields", "набор полей формы", "один набор fields = один экран в приложении", ACCENT),
        ("flow", "последовательность экранов", "flow = сценарий сбора (набор экранов по порядку)", ACCENT_PURPLE),
        ("viz.json", "отображение в админке", "визуализации и layout просмотра пакета", ACCENT_GREEN),
    ]
    for title, sub, hint, c in parts:
        _box(slide, ix, y, bw, bh, BG_CARD, c, title, sub, 11, 8)
        _caption(slide, ix + Inches(0.08), y + bh + Inches(0.02), bw - Inches(0.16), hint, TEXT_MUTED, 7)
        y += bh + hint_h + Inches(0.1)
        if title != parts[-1][0]:
            _arrow_v(slide, ix + bw / 2, y - Inches(0.1), y - Inches(0.02), c)

    foot_y = t + h - Inches(0.24)
    _caption(slide, ix, foot_y, iw,
             "Новый сценарий = правка JSON в Git, без релиза приложения", ACCENT_CYAN, 8)


# ── Package viewing (Step 4) ──────────────────────────────────────────────────

def draw_package_viewing_flow(slide, l, t, w, h):
    """List → open package → workspace tabs."""
    _canvas(slide, l, t, w, h)
    pad = Inches(0.14)
    ix, iy = l + pad, t + Inches(0.2)
    iw = w - pad * 2
    cx = ix + iw / 2

    # 1 List
    lh = Inches(0.68)
    _box(slide, ix + iw * 0.15, iy, iw * 0.7, lh, BG_CARD, ACCENT_GREEN,
         "/ui/packages/", "список принятых пакетов", 11, 8)

    # 2 filters row
    fy = iy + lh + Inches(0.32)
    fw = (iw - Inches(0.2)) / 3
    for i, label in enumerate(["проект", "статус", "поиск"]):
        _box(slide, ix + Inches(0.1) + i * (fw + Inches(0.05)), fy, fw, Inches(0.42),
             BG_PANEL, BORDER, label, "", 9, 7)
    _arrow_v(slide, cx, iy + lh, fy, ACCENT_GREEN)

    # 3 workspace
    wy = fy + Inches(0.58)
    wh = Inches(0.62)
    _box(slide, ix + iw * 0.1, wy, iw * 0.8, wh, RGBColor(0x1A, 0x2E, 0x28), ACCENT_GREEN,
         "Package Workspace", "один пакет · все данные", 11, 8)
    _arrow_v(slide, cx, fy + Inches(0.42), wy, ACCENT_GREEN)

    # 4 tabs
    ty = wy + wh + Inches(0.28)
    tw = (iw - Inches(0.09 * 3)) / 4
    th = t + h - ty - Inches(0.12)
    tabs = [
        ("Данные", ACCENT),
        ("Медиа", ACCENT_PURPLE),
        ("Viz", ACCENT_GREEN),
        ("История", ACCENT_GOLD),
    ]
    for i, (tab, c) in enumerate(tabs):
        _box(slide, ix + i * (tw + Inches(0.09)), ty, tw, th, BG_CARD, c, tab, "", 10, 7)
    _arrow_v(slide, cx, wy + wh, ty, ACCENT_GREEN)


def draw_admin_workspace_mini(slide, l, t, w, h):
    """Compact 4-tab view for step slide sidebar."""
    draw_package_viewing_flow(slide, l, t, w, h)


# Legacy alias for client-server — now same as five step
def draw_client_server_flow(slide, l, t, w, h):
    draw_five_step_flow(slide, l, t, w, h)


def draw_admin_workspace(slide, l, t, w, h):
    draw_package_viewing_flow(slide, l, t, w, h)
