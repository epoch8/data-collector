"""Generate Data Collector presentation (v3 — clean diagrams, lean deck)."""

from pathlib import Path

from diagrams import (
    draw_config_entities,
    draw_cover_architecture,
    draw_five_step_flow,
    draw_package_viewing_flow,
    draw_stack_diagram,
)
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

BG_DEEP = RGBColor(0x0A, 0x0E, 0x18)
BG = RGBColor(0x12, 0x12, 0x18)
BG_ELEVATED = RGBColor(0x18, 0x20, 0x30)
BG_CARD = RGBColor(0x1E, 0x2A, 0x3D)
ACCENT = RGBColor(0x5B, 0x8D, 0xEF)
ACCENT_SOFT = RGBColor(0x3A, 0x5A, 0x9E)
ACCENT_GREEN = RGBColor(0x3D, 0xCC, 0x85)
ACCENT_GOLD = RGBColor(0xD4, 0xA8, 0x4B)
ACCENT_PURPLE = RGBColor(0x9B, 0x7E, 0xD0)
ACCENT_CORAL = RGBColor(0xE8, 0x7A, 0x7E)
TEXT = RGBColor(0xE8, 0xEA, 0xED)
TEXT_MUTED = RGBColor(0x9C, 0xA3, 0xAF)
TEXT_DIM = RGBColor(0x6B, 0x72, 0x80)
BORDER = RGBColor(0x3A, 0x45, 0x58)
PLACEHOLDER_BG = RGBColor(0x1A, 0x20, 0x2C)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)
TOTAL_SLIDES = 11


def _gradient_bg(slide, c0=BG_DEEP, c1=BG_ELEVATED, angle=135):
    fill = slide.background.fill
    fill.gradient()
    fill.gradient_angle = angle
    fill.gradient_stops[0].color.rgb = c0
    fill.gradient_stops[1].color.rgb = c1


def _shape(slide, kind, left, top, w, h, fill, line=None, line_w=Pt(1)):
    s = slide.shapes.add_shape(kind, left, top, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = fill
    if line:
        s.line.color.rgb = line
        s.line.width = line_w
    else:
        s.line.fill.background()
    return s


def _rect(slide, l, t, w, h, fill, line=None):
    return _shape(slide, MSO_SHAPE.RECTANGLE, l, t, w, h, fill, line)


def _round(slide, l, t, w, h, fill, line=None):
    return _shape(slide, MSO_SHAPE.ROUNDED_RECTANGLE, l, t, w, h, fill, line)


def _circle(slide, l, t, size, fill):
    return _shape(slide, MSO_SHAPE.OVAL, l, t, size, size, fill)


def _text(slide, l, t, w, h, text, size=16, bold=False, color=TEXT, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    box = slide.shapes.add_textbox(l, t, w, h)
    tf = box.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.alignment = align
    return box


def _decor_orbs(slide):
    _circle(slide, Inches(-1.2), Inches(-1.0), Inches(3.5), ACCENT_SOFT)
    _circle(slide, Inches(11.5), Inches(5.8), Inches(2.8), RGBColor(0x2A, 0x35, 0x50))


def _slide_chrome(slide, num, eyebrow=None, accent=ACCENT):
    _gradient_bg(slide)
    _decor_orbs(slide)
    _rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.06), accent)
    _rect(slide, Inches(0), Inches(7.18), SLIDE_W, Inches(0.32), BG_DEEP)
    _text(slide, Inches(0.65), Inches(7.2), Inches(4), Inches(0.28), "Data Collector  ·  Epoch8", 9, color=TEXT_DIM)
    _text(slide, Inches(12.0), Inches(7.2), Inches(0.9), Inches(0.28), f"{num:02d} / {TOTAL_SLIDES}", 9, color=TEXT_DIM, align=PP_ALIGN.RIGHT)
    if eyebrow:
        pill = _round(slide, Inches(0.65), Inches(0.45), Inches(1.8), Inches(0.32), BG_CARD, accent)
        pill.text_frame.paragraphs[0].text = eyebrow
        pill.text_frame.paragraphs[0].font.size = Pt(9)
        pill.text_frame.paragraphs[0].font.bold = True
        pill.text_frame.paragraphs[0].font.color.rgb = accent
        pill.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        pill.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE


def _slide_title(slide, title, subtitle=None, top=Inches(0.85)):
    _text(slide, Inches(0.65), top, Inches(11), Inches(0.75), title, 30, bold=True)
    _rect(slide, Inches(0.65), top + Inches(0.78), Inches(2.2), Inches(0.05), ACCENT)
    if subtitle:
        _text(slide, Inches(0.65), top + Inches(0.95), Inches(11), Inches(0.45), subtitle, 13, color=TEXT_MUTED)


def _screenshot_frame(slide, l, t, w, h, label, accent=ACCENT):
    outer = _round(slide, l, t, w, h, PLACEHOLDER_BG, accent)
    outer.text_frame.clear()
    pad = Inches(0.1)
    inner = _round(slide, l + pad, t + pad, w - pad * 2, h - pad * 2, BG_ELEVATED, BORDER)
    inner.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    p = inner.text_frame.paragraphs[0]
    p.text = label
    p.font.size = Pt(10)
    p.font.color.rgb = TEXT_MUTED
    p.alignment = PP_ALIGN.CENTER
    tag = _round(slide, l + w - Inches(1.1), t + Inches(0.08), Inches(1.0), Inches(0.24), accent)
    tag.text_frame.paragraphs[0].text = "SCREENSHOT"
    tag.text_frame.paragraphs[0].font.size = Pt(6)
    tag.text_frame.paragraphs[0].font.bold = True
    tag.text_frame.paragraphs[0].font.color.rgb = TEXT
    tag.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    tag.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE


def _bullet_cards(slide, items, l, t, w, accent=ACCENT):
    y = t
    for item in items:
        if not item.strip():
            y += Inches(0.1)
            continue
        is_solution = item.lower().startswith("решение")
        is_problem = item.lower().startswith("проблема")
        card_fill = BG_CARD if not is_solution else RGBColor(0x14, 0x28, 0x22)
        card_line = ACCENT_CORAL if is_problem else (ACCENT_GREEN if is_solution else accent)
        card_h = Inches(0.58) if len(item) < 72 else Inches(0.76)
        _rect(slide, l, y + Inches(0.06), Inches(0.05), card_h - Inches(0.12), card_line)
        card = _round(slide, l + Inches(0.08), y, w - Inches(0.08), card_h, card_fill, BORDER)
        card.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = card.text_frame.paragraphs[0]
        p.text = item
        p.font.size = Pt(13)
        p.font.color.rgb = TEXT if not is_solution else ACCENT_GREEN
        p.font.bold = is_solution
        y += card_h + Inches(0.08)


def _component_cards(slide, rows, l, t, w):
    icons = ["⚙", "👁", "📱", "☁"]
    colors = [ACCENT_GOLD, ACCENT_GREEN, ACCENT_PURPLE, ACCENT]
    card_h = Inches(1.05)
    gap = Inches(0.1)
    for i, (name, who, what) in enumerate(rows):
        y = t + i * (card_h + gap)
        c = colors[i % len(colors)]
        card = _round(slide, l, y, w, card_h, BG_CARD, BORDER)
        card.text_frame.clear()
        badge = _round(slide, l + Inches(0.12), y + Inches(0.2), Inches(0.5), Inches(0.5), RGBColor(0x28, 0x32, 0x48), c)
        badge.text_frame.paragraphs[0].text = icons[i]
        badge.text_frame.paragraphs[0].font.size = Pt(14)
        badge.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        badge.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        _text(slide, l + Inches(0.75), y + Inches(0.1), Inches(2), Inches(0.32), name, 12, bold=True, color=c)
        _text(slide, l + Inches(0.75), y + Inches(0.38), Inches(2), Inches(0.28), who, 9, color=TEXT_DIM)
        _text(slide, l + Inches(2.85), y + Inches(0.2), w - Inches(3), Inches(0.65), what, 11, color=TEXT_MUTED)


def _step_badge(slide, num, color, l, t):
    b = _circle(slide, l, t, Inches(0.5), color)
    b.text_frame.paragraphs[0].text = str(num)
    b.text_frame.paragraphs[0].font.size = Pt(15)
    b.text_frame.paragraphs[0].font.bold = True
    b.text_frame.paragraphs[0].font.color.rgb = BG
    b.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    b.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE


def _slide_step(prs, num, title, eyebrow, bullets, step_num, step_color, screenshot_label=None, diagram_fn=None, subtitle=None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_chrome(slide, num, eyebrow, step_color)
    title_top = Inches(1.02)
    _step_badge(slide, step_num, step_color, Inches(0.65), Inches(0.38))
    _slide_title(slide, title, subtitle, title_top)
    _bullet_cards(slide, bullets, Inches(0.65), title_top + Inches(1.0), Inches(6.2), step_color)
    rl, rt, rw, rh = Inches(7.15), Inches(1.05), Inches(5.55), Inches(5.85)
    if diagram_fn:
        diagram_fn(slide, rl, rt, rw, rh)
    elif screenshot_label:
        _screenshot_frame(slide, rl, rt, rw, rh, screenshot_label, step_color)
    return slide


def build() -> Path:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # 1 — Cover
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _gradient_bg(slide, BG_DEEP, RGBColor(0x14, 0x1C, 0x32), 145)
    _decor_orbs(slide)
    _rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), ACCENT)
    panel = _round(slide, Inches(0.55), Inches(0.75), Inches(7.2), Inches(6.0), BG_CARD, BORDER)
    _text(slide, Inches(0.95), Inches(1.05), Inches(2), Inches(0.3), "EPOCH8", 10, bold=True, color=ACCENT_GOLD)
    _text(slide, Inches(0.95), Inches(1.55), Inches(6.5), Inches(1.1), "Data Collector", 44, bold=True)
    _rect(slide, Inches(0.95), Inches(2.75), Inches(2.4), Inches(0.06), ACCENT)
    _text(slide, Inches(0.95), Inches(3.0), Inches(6.4), Inches(1.4),
          "Платформа сбора полевых данных\nОдин бэкенд · мобильное приложение · конфиг без релиза", 17, color=TEXT_MUTED)
    draw_cover_architecture(slide, Inches(8.1), Inches(0.75), Inches(4.65), Inches(2.8))
    _screenshot_frame(slide, Inches(8.1), Inches(3.7), Inches(4.65), Inches(3.05), "hero — админка + телефон")

    # 2 — Overview
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_chrome(slide, 2, "OVERVIEW")
    _slide_title(slide, "Зачем это нужно", "От хаоса в поле — к структурированным данным")
    _bullet_cards(slide, [
        "Проблема: каждый сценарий сбора = новая форма в коде",
        "Проблема: фото и метаданные теряются в мессенджерах и Excel",
        "Проблема: нет единого места «что собрали» и «что дальше»",
        "",
        "Решение: конфиг → UI в приложении → пакет на сервер → просмотр в админке",
    ], Inches(0.65), Inches(1.95), Inches(6.2))
    _screenshot_frame(slide, Inches(7.15), Inches(1.05), Inches(5.55), Inches(5.85), "до / после — хаос vs список пакетов")

    # 3 — Product
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_chrome(slide, 3, "PRODUCT")
    _slide_title(slide, "Из чего состоит продукт", "Одна платформа — N проектов")
    _component_cards(slide, [
        ("Staff Admin", "Epoch8", "конфиг · Git · пользователи"),
        ("Client Admin", "Заказчик", "просмотр пакетов · viz"),
        ("Flutter App", "Сборщик", "сбор · офлайн · upload"),
        ("Django + storage", "Платформа", "API · БД · медиа"),
    ], Inches(0.65), Inches(1.95), Inches(6.2))
    draw_stack_diagram(slide, Inches(7.15), Inches(1.05), Inches(5.55), Inches(5.85))

    # 4 — Flow
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_chrome(slide, 4, "FLOW")
    _slide_title(slide, "Путь данных", "5 шагов от настройки до просмотра")
    draw_five_step_flow(slide, Inches(0.65), Inches(2.15), Inches(12.05), Inches(4.35))
    _text(
        slide, Inches(0.65), Inches(6.65), Inches(12.05), Inches(0.35),
        "Один конфиг → тот же сценарий в мобилке → пакет той же структуры в админке",
        12, color=TEXT_MUTED, align=PP_ALIGN.CENTER,
    )

    # 5 — Config
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_chrome(slide, 5, "CONFIG")
    _slide_title(slide, "Конфиг = сценарий без кода", "JSON в Git описывает сбор и просмотр")
    _bullet_cards(slide, [
        "fields — какие поля собираем (текст, фото, дата…)",
        "flow — порядок экранов в мобилке",
        "viz.json — как показывать пакет в админке",
        "",
        "Новый сценарий = правка JSON, не релиз приложения",
    ], Inches(0.65), Inches(1.95), Inches(6.2))
    draw_config_entities(slide, Inches(7.15), Inches(1.95), Inches(5.55), Inches(2.2))
    _screenshot_frame(slide, Inches(7.15), Inches(4.35), Inches(5.55), Inches(2.55), "JSON-редактор (опционально)")

    # 6–8 — Steps 1–3
    _slide_step(prs, 6, "Админ настраивает проект", "STEP 01", [
        "Создаёт проект и привязывает Git",
        "Редактирует config.json: fields + flow",
        "Назначает сборщиков на проект",
        "Сохраняет → commit + push",
    ], 1, ACCENT_GOLD, screenshot_label="проекты · JSON-редактор · пользователи")

    _slide_step(prs, 7, "Конфиг в приложении", "STEP 02", [
        "Сборщик логинится и видит свои проекты",
        "Приложение скачивает config с сервера",
        "Кэширует локально — работает офлайн",
        "Разные конфиги → разный UI в одном app",
    ], 2, ACCENT, screenshot_label="список проектов · форма · камера")

    _slide_step(prs, 8, "Сборщик сдаёт пакет", "STEP 03", [
        "Проходит сценарий по config.flow",
        "Пакет = поля + фото + манифест",
        "Черновик сохраняется локально",
        "Upload на сервер при появлении сети",
    ], 3, ACCENT_PURPLE, screenshot_label="форма · review · статус загрузки")

    # 9 — Step 4 intro
    _slide_step(prs, 9, "Пакет принят на сервер", "STEP 04", [
        "После upload пакет появляется в админке",
        "Привязан к project_id и config версии",
        "Манифест + blobs хранятся на сервере",
        "Готов к просмотру и анализу",
    ], 4, ACCENT_GREEN, diagram_fn=draw_package_viewing_flow)

    # 10 — Step 4 deep dive: viewing packages
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_chrome(slide, 10, "STEP 04", ACCENT_GREEN)
    _step_badge(slide, 4, ACCENT_GREEN, Inches(0.65), Inches(0.38))
    _slide_title(slide, "Просмотр принятых пакетов", "Админка: список → фильтр → workspace пакета")
    _bullet_cards(slide, [
        "Список /ui/packages/ — все принятые пакеты проекта",
        "Фильтры: проект, статус, поиск по полям",
        "Workspace: данные · медиа · viz · история правок",
        "UI строится из тех же fields, что и при сборе",
    ], Inches(0.65), Inches(2.05), Inches(5.8), ACCENT_GREEN)
    draw_package_viewing_flow(slide, Inches(6.75), Inches(1.05), Inches(6.0), Inches(3.55))
    _screenshot_frame(slide, Inches(6.75), Inches(4.75), Inches(2.9), Inches(2.15), "список пакетов", ACCENT_GREEN)
    _screenshot_frame(slide, Inches(9.85), Inches(4.75), Inches(2.9), Inches(2.15), "workspace · вкладки", ACCENT_GREEN)

    # 11 — Closing
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _gradient_bg(slide, RGBColor(0x0E, 0x18, 0x28), BG_DEEP, 160)
    _decor_orbs(slide)
    _rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), ACCENT_GREEN)
    _round(slide, Inches(1.2), Inches(1.5), Inches(10.9), Inches(4.5), BG_CARD, ACCENT_GREEN).text_frame.clear()
    _text(slide, Inches(1.65), Inches(1.85), Inches(10), Inches(0.6), "Data Collector в одном предложении", 28, bold=True, align=PP_ALIGN.CENTER)
    _text(slide, Inches(1.65), Inches(2.65), Inches(10), Inches(2.2),
          "Админ описывает сценарий в JSON → приложение проводит сборщика по шагам → "
          "пакет с медиа попадает в админку → команда видит данные без ручной склейки.",
          18, color=TEXT_MUTED, align=PP_ALIGN.CENTER)
    cta = _round(slide, Inches(3.5), Inches(5.0), Inches(6.3), Inches(0.55), ACCENT)
    cta.text_frame.paragraphs[0].text = "Опишите поля и flow — соберём demo под ваш кейс"
    cta.text_frame.paragraphs[0].font.size = Pt(13)
    cta.text_frame.paragraphs[0].font.bold = True
    cta.text_frame.paragraphs[0].font.color.rgb = TEXT
    cta.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    cta.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
    _text(slide, Inches(12.0), Inches(7.2), Inches(0.9), Inches(0.28), f"{TOTAL_SLIDES:02d} / {TOTAL_SLIDES}", 9, color=TEXT_DIM, align=PP_ALIGN.RIGHT)

    out = Path(__file__).resolve().parent / "Data-Collector-Canva-Template-v3.pptx"
    prs.save(out)
    return out


if __name__ == "__main__":
    print(f"Created: {build()}")
