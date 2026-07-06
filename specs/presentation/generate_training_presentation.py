"""Generate training program presentation (plan + implementation roadmap)."""

from pathlib import Path

from diagrams import (
    draw_eight_blocks_timeline,
    draw_learning_formats,
    draw_parallel_tracks,
    draw_stand_projects,
    draw_training_journey,
)
from generate_canva_template import (
    ACCENT,
    ACCENT_GOLD,
    ACCENT_GREEN,
    ACCENT_PURPLE,
    BG_CARD,
    BG_DEEP,
    BORDER,
    SLIDE_H,
    SLIDE_W,
    TEXT,
    TEXT_DIM,
    TEXT_MUTED,
    _bullet_cards,
    _circle,
    _component_cards,
    _decor_orbs,
    _gradient_bg,
    _rect,
    _round,
    _slide_title,
    _text,
)
import generate_canva_template as slide_template
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

TOTAL_SLIDES = 14
OUT_PATH = Path(__file__).resolve().parent / "Training-Program.pptx"

slide_template.TOTAL_SLIDES = TOTAL_SLIDES
_slide_chrome = slide_template._slide_chrome

ACCENT_CORAL = RGBColor(0xE8, 0x7A, 0x7E)
ACCENT_CYAN = RGBColor(0x8B, 0xC9, 0xF7)


def _stat_card(slide, l, t, w, h, value, label, accent=ACCENT):
    card = _round(slide, l, t, w, h, BG_CARD, accent)
    card.text_frame.clear()
    p0 = card.text_frame.paragraphs[0]
    p0.text = value
    p0.font.size = Pt(28)
    p0.font.bold = True
    p0.font.color.rgb = accent
    p0.alignment = PP_ALIGN.CENTER
    p1 = card.text_frame.add_paragraph()
    p1.text = label
    p1.font.size = Pt(10)
    p1.font.color.rgb = TEXT_MUTED
    p1.alignment = PP_ALIGN.CENTER
    card.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE


def _mini_table_rows(slide, rows, l, t, w, accent=ACCENT):
    y = t
    for i, (c1, c2) in enumerate(rows):
        color = accent if i == 0 else BORDER
        row_h = Inches(0.42) if len(c2) < 48 else Inches(0.56)
        card = _round(slide, l, y, w, row_h, BG_CARD if i else RGBColor(0x1A, 0x24, 0x36), color)
        card.text_frame.clear()
        p0 = card.text_frame.paragraphs[0]
        p0.text = c1
        p0.font.size = Pt(10 if i else 9)
        p0.font.bold = bool(i == 0)
        p0.font.color.rgb = accent if i == 0 else TEXT
        if c2:
            p1 = card.text_frame.add_paragraph()
            p1.text = c2
            p1.font.size = Pt(9)
            p1.font.color.rgb = TEXT_MUTED
        card.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        y += row_h + Inches(0.06)


def build() -> Path:
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    # 1 — Cover
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _gradient_bg(slide, BG_DEEP, RGBColor(0x14, 0x1C, 0x32), 145)
    _decor_orbs(slide)
    _rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), ACCENT_GREEN)
    panel = _round(slide, Inches(0.55), Inches(0.75), Inches(7.35), Inches(6.0), BG_CARD, BORDER)
    _text(slide, Inches(0.95), Inches(1.05), Inches(2.5), Inches(0.3), "E8 team  ·  Обучение", 10, bold=True, color=ACCENT_GOLD)
    _text(slide, Inches(0.95), Inches(1.55), Inches(6.8), Inches(1.35),
          "Программа обучения\nкоманды сопровождения", 40, bold=True)
    _rect(slide, Inches(0.95), Inches(3.05), Inches(2.6), Inches(0.06), ACCENT_GREEN)
    _text(slide, Inches(0.95), Inches(3.35), Inches(6.5), Inches(1.6),
          "План 300 ч + дорожная карта реализации\n"
          "data-collector · datapipe · бонитировка КРС", 17, color=TEXT_MUTED)
    _text(slide, Inches(0.95), Inches(5.35), Inches(6.2), Inches(0.8),
          "Full-stack Engineer  +  ML Engineer", 13, color=ACCENT_CYAN)
    draw_training_journey(slide, Inches(8.15), Inches(1.05), Inches(4.65), Inches(5.55))

    # 2 — Program at a glance
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_chrome(slide, 2, "OVERVIEW", ACCENT_GREEN)
    _slide_title(slide, "Программа в цифрах", "Утвержденный план и формат реализации")
    cw = Inches(2.85)
    ch = Inches(1.35)
    gap = Inches(0.18)
    stats = [
        ("300 ч", "работа e8", ACCENT_GREEN),
        ("990 000 ₽", "бюджет программы", ACCENT_GOLD),
        ("15 нед.", "календарь", ACCENT),
        ("8 блоков", "структура курса", ACCENT_PURPLE),
        ("2 роли", "Full-stack + ML", ACCENT_CYAN),
        ("~394 ч", "освоение на участника", ACCENT_CORAL),
    ]
    for i, (val, lbl, col) in enumerate(stats):
        row, col_i = divmod(i, 3)
        _stat_card(slide, Inches(0.65) + col_i * (cw + gap), Inches(2.05) + row * (ch + gap), cw, ch, val, lbl, col)
    _bullet_cards(slide, [
        "Цель: не лекции, а проектное погружение в реальный продуктовый цикл",
        "Материалы делятся на Reusable Core, Product Layer и Korovas Layer",
        "Каждый блок заканчивается проверяемым артефактом: отчет, PR, E2E, runbook",
    ], Inches(0.65), Inches(5.0), Inches(12.0), ACCENT_GREEN)

    # 3 — Philosophy
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_chrome(slide, 3, "APPROACH", ACCENT)
    _slide_title(slide, "Как устроено обучение", "Сценарии вместо теории ради теории")
    _bullet_cards(slide, [
        "Участник проходит путь от предметной области до аттестации на стенде",
        "Короткие видео снимают повторяющееся объяснение — синхронное время для разборов",
        "Практикумы и разборы инцидентов дают ощущение реальной эксплуатации",
        "Ревью кода и пайплайна проверяет способность вести проект самостоятельно",
        "Финал: перенос решений в новую предметную область — например, бонитировка баранов",
    ], Inches(0.65), Inches(1.95), Inches(6.2))
    draw_training_journey(slide, Inches(7.15), Inches(1.05), Inches(5.55), Inches(5.85))

    # 4 — Learning formats
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_chrome(slide, 4, "FORMAT", ACCENT_PURPLE)
    _slide_title(slide, "Форматы обучения", "Смешанная программа: видео + live + практика")
    draw_learning_formats(slide, Inches(0.65), Inches(1.85), Inches(6.35), Inches(4.85))
    _bullet_cards(slide, [
        "Микровидео: 7-15 минут, один конкретный вопрос",
        "Живой семинар: демо, Q&A, совместный разбор пакета",
        "Практикум: участник делает руками и видит результат в системе",
        "Разбор инцидента: заранее сломанные upload, datapipe, viz, storage",
        "Ролевое задание и ревью: проверка самостоятельной работы",
    ], Inches(7.25), Inches(1.95), Inches(5.45), ACCENT_PURPLE)

    # 5 — Eight blocks
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_chrome(slide, 5, "STRUCTURE", ACCENT_GOLD)
    _slide_title(slide, "8 блоков программы", "От предметки до углубления и аттестации")
    draw_eight_blocks_timeline(slide, Inches(0.65), Inches(1.85), Inches(12.05), Inches(4.95))

    # 6 — Blocks detail 1-4
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_chrome(slide, 6, "BLOCKS", ACCENT)
    _slide_title(slide, "Блоки 1-4", "Онбординг, инженерия, архитектура, full-stack")
    rows = [
        ("Блок", "Содержание и проверка"),
        ("1. Предметная область", "Бонитировка, CVAT, путь пакета · квиз"),
        ("2. Инженерная база", "Git/PR, API, Cursor, локальный стенд · PR + техотчет"),
        ("3. Архитектура", "Компоненты, поток данных, E2E · схема + отчет"),
        ("4. Full-stack", "Flutter, Django, upload, viz · E2E-сценарий + ревью"),
    ]
    _mini_table_rows(slide, rows, Inches(0.65), Inches(1.95), Inches(6.1))
    _bullet_cards(slide, [
        "Блоки 1-3 общие для обеих ролей",
        "Фокус: понять продукт, архитектуру и первый сквозной сценарий",
        "Проверка через квиз, ревью и пошаговые отчеты со скриншотами",
    ], Inches(7.15), Inches(1.95), Inches(5.55))

    # 7 — Blocks detail 5-8
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_chrome(slide, 7, "BLOCKS", ACCENT_GREEN)
    _slide_title(slide, "Блоки 5-8", "ML, эксплуатация, прод, углубление")
    rows = [
        ("Блок", "Содержание и проверка"),
        ("5. ML + datapipe", "Pipeline, CVAT, inference, метрики · запуск + отчет"),
        ("6. Эксплуатация", "Деплой, логи, инциденты, runbook · отчет + ревью"),
        ("7. Практика в проде", "3 задания из backlog · ревью + устная защита"),
        ("8. Углубление", "Новая предметная область · PR + документация"),
    ]
    _mini_table_rows(slide, rows, Inches(0.65), Inches(1.95), Inches(6.1))
    draw_parallel_tracks(slide, Inches(7.15), Inches(1.85), Inches(5.55), Inches(5.05))

    # 8 — Data Collector materials
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_chrome(slide, 8, "MATERIALS", ACCENT_PURPLE)
    _slide_title(slide, "Материалы: Data Collector", "Сбор, upload, админка, визуализация")
    _component_cards(slide, [
        ("Flutter-клиент", "Сборщик", "wizard · Drift · история · sync"),
        ("Конфиг проекта", "Staff", "collector/config.json · сценарий без релиза"),
        ("Пакет и upload", "Платформа", "payload.json · blobs · POST/PUT/commit"),
        ("Админка + viz", "Заказчик", "пакеты · медиа · collector/viz.json"),
    ], Inches(0.65), Inches(1.95), Inches(6.2))
    _bullet_cards(slide, [
        "Практикумы: изменить конфиг, разобрать payload, пройти upload protocol",
        "Разбор инцидентов: пустая визуализация, сбой upload, роли и доступы",
        "Материалы переиспользуются для будущих заказчиков без доменной привязки",
    ], Inches(7.15), Inches(1.95), Inches(5.55), ACCENT_PURPLE)

    # 9 — DataPipe materials
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_chrome(slide, 9, "MATERIALS", ACCENT)
    _slide_title(slide, "Материалы: DataPipe", "Пайплайны, CVAT, inference, метрики")
    _component_cards(slide, [
        ("Карта datapipe", "Обе роли", "вход · шаг · артефакт · результат"),
        ("Project DB", "ML", "annotations · inference · depth · CVAT links"),
        ("CVAT round-trip", "ML", "экспорт/импорт разметки по пакету"),
        ("Качество модели", "ML", "PCK · версии · дообучение · rollback"),
    ], Inches(0.65), Inches(1.95), Inches(6.2))
    _bullet_cards(slide, [
        "Отдельный пакет материалов, но те же учебные данные что и в data-collector",
        "Практикум: запустить inference до появления результата в админке",
        "Отчет по метрикам: сравнение GT, модели и эксперта",
    ], Inches(7.15), Inches(1.95), Inches(5.55))

    # 10 — Korovas layer
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_chrome(slide, 10, "DOMAIN", ACCENT_GOLD)
    _slide_title(slide, "Korovas Layer", "Доменная специфика текущего заказчика")
    _bullet_cards(slide, [
        "21 точка · 15 параметров · 9 признаков бонитировки КРС",
        "Правила съемки: интерактивный чек-лист пригодных фото для inference",
        "Галерея кейсов: хорошие и плохие примеры съемки и разметки",
        "Семинар «эксперт vs модель»: объяснение расхождений результата",
        "Набор пакетов по коровам — основа практикумов, инцидентов и аттестации",
    ], Inches(0.65), Inches(1.95), Inches(6.2), ACCENT_GOLD)
    card = _round(slide, Inches(7.15), Inches(1.95), Inches(5.55), Inches(5.05), BG_CARD, ACCENT_GOLD)
    card.text_frame.clear()
    lines = [
        ("Reusable Core", "Git · API · архитектура · эксплуатация"),
        ("Product Layer", "data-collector + datapipe"),
        ("Korovas Layer", "бонитировка · съемка · метрики КРС"),
    ]
    y = Inches(2.35)
    for title, sub in lines:
        _text(slide, Inches(7.55), y, Inches(4.8), Inches(0.35), title, 14, bold=True, color=ACCENT_GOLD)
        _text(slide, Inches(7.55), y + Inches(0.38), Inches(4.8), Inches(0.35), sub, 11, color=TEXT_MUTED)
        y += Inches(1.15)

    # 11 — Training stand
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_chrome(slide, 11, "STAND", ACCENT_GREEN)
    _slide_title(slide, "Учебный стенд", "Безопасная среда для практики и инцидентов")
    draw_stand_projects(slide, Inches(0.65), Inches(1.85), Inches(5.85), Inches(5.05))
    _bullet_cards(slide, [
        "mobile / emulator → staging API → project DB + storage",
        "datapipe staging → CVAT → визуализация в админке",
        "Frozen packages для повторяемых labs, mutable — для заданий",
        "korovas-broken: набор намеренно сломанных кейсов для блока эксплуатации",
        "Reset-инструкция для менторов перед каждой неделей",
    ], Inches(6.75), Inches(1.95), Inches(5.95), ACCENT_GREEN)

    # 12 — E2E scenarios
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_chrome(slide, 12, "SCENARIOS", ACCENT_CYAN)
    _slide_title(slide, "Сценарии на стенде", "Что участник должен уметь воспроизвести")
    scenarios = [
        ("Снять пакет", "локально готов к upload"),
        ("Загрузить пакет", "session · blobs · manifest · completed"),
        ("Проверить хранилище", "blob + запись в project DB"),
        ("Запустить datapipe", "inference / depth / CVAT rows"),
        ("Открыть визуализацию", "слои GT · inference · CVAT · depth"),
        ("Разобрать сбой", "гипотеза · логи · фикс · эскалация"),
        ("Сравнить модель и эксперта", "отчет по метрикам"),
    ]
    y = Inches(1.95)
    for i, (title, result) in enumerate(scenarios):
        colors = [ACCENT, ACCENT_PURPLE, ACCENT_GREEN, ACCENT_GOLD, ACCENT_CYAN, ACCENT_CORAL, ACCENT]
        c = colors[i % len(colors)]
        badge = _circle(slide, Inches(0.75), y + Inches(0.04), Inches(0.38), c)
        badge.text_frame.paragraphs[0].text = str(i + 1)
        badge.text_frame.paragraphs[0].font.size = Pt(11)
        badge.text_frame.paragraphs[0].font.bold = True
        badge.text_frame.paragraphs[0].font.color.rgb = BG_DEEP
        badge.text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        badge.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
        _text(slide, Inches(1.25), y, Inches(3.2), Inches(0.32), title, 12, bold=True, color=c)
        _text(slide, Inches(4.5), y, Inches(8.2), Inches(0.32), result, 11, color=TEXT_MUTED)
        y += Inches(0.62)

    # 13 — Assessment
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _slide_chrome(slide, 13, "ASSESSMENT", ACCENT_CORAL)
    _slide_title(slide, "Проверка и аттестация", "Артефакты, ревью, устная защита")
    _component_cards(slide, [
        ("Квиз", "Блок 1", "предметная область"),
        ("Ревью", "Блоки 2-8", "PR · pipeline · отчеты"),
        ("Техотчет", "Все блоки", "симптомы · логи · фикс · проверка"),
        ("Аттестация", "Блок 7", "3 задания + устная защита"),
    ], Inches(0.65), Inches(1.95), Inches(6.2))
    _bullet_cards(slide, [
        "Full-stack: E2E, upload, viz, деплой, роли и доступы",
        "ML: pipeline run, CVAT, метрики, сравнение моделей, план дообучения",
        "Финал: перенос решений в смежную предметную область",
        "Критерий: участник объясняет свой участок и границы соседних систем",
    ], Inches(7.15), Inches(1.95), Inches(5.55), ACCENT_CORAL)

    # 14 — Closing
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _gradient_bg(slide, BG_DEEP, RGBColor(0x14, 0x28, 0x22), 135)
    _decor_orbs(slide)
    _rect(slide, Inches(0), Inches(0), SLIDE_W, Inches(0.08), ACCENT_GREEN)
    _text(slide, Inches(0.95), Inches(1.35), Inches(11), Inches(0.9), "Что получаем на выходе", 34, bold=True)
    _rect(slide, Inches(0.95), Inches(2.35), Inches(2.4), Inches(0.06), ACCENT_GREEN)
    _bullet_cards(slide, [
        "Команда заказчика умеет вести data-collector и datapipe end-to-end",
        "Готовый пакет reusable-материалов для следующих клиентов",
        "Runbooks, шаблоны техотчетов и учебный стенд для повторного запуска",
        "Понятные критерии приемки по ролям Full-stack и ML Engineer",
    ], Inches(0.95), Inches(2.65), Inches(7.5), ACCENT_GREEN)
    _text(slide, Inches(0.95), Inches(6.35), Inches(8), Inches(0.4),
          "Epoch8  ·  Training Program  ·  2026", 11, color=TEXT_DIM)
    _text(slide, Inches(11.2), Inches(7.2), Inches(1.5), Inches(0.28), f"{TOTAL_SLIDES:02d} / {TOTAL_SLIDES}", 9, color=TEXT_DIM, align=PP_ALIGN.RIGHT)

    prs.save(OUT_PATH)
    return OUT_PATH


if __name__ == "__main__":
    path = build()
    print(f"Saved: {path}")
